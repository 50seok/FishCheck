"""
방어 vs 부시리 분류 테스트 + PPT 생성 스크립트.

사용법:
  python tools/test_amberjack_ppt.py

입력 폴더:
  data/raw/bangeo/   → 예상 클래스: bangeo
  data/raw/bushiri/  → 예상 클래스: bushiri

출력: ppt/방어_부시리_분류테스트.pptx
"""
import os
import sys
import io

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)

from PIL import Image
from ultralytics import YOLO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

MODEL_PATH  = os.path.join(ROOT, "models", "best.pt")
OUTPUT_PATH = os.path.join(ROOT, "ppt", "방어_부시리_분류테스트_v2.pptx")

TEST_SETS = {
    "bangeo":  os.path.join(ROOT, "data", "raw", "bangeo"),
    "bushiri": os.path.join(ROOT, "data", "raw", "bushiri"),
}

CLASS_KO = {
    "bangeo":  "방어",
    "bushiri": "부시리",
    "gajami":  "가자미/도다리",
    "gwangeo": "광어",
    "none":    "미탐지",
}

IMG_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".bmp"}

GREEN  = RGBColor(0x2E, 0x86, 0x48)
RED    = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BLUE   = RGBColor(0x1A, 0x5C, 0x96)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF7, 0xFA)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def collect_images(folder):
    return [
        os.path.join(folder, fn)
        for fn in sorted(os.listdir(folder))
        if os.path.splitext(fn)[1].lower() in IMG_EXTS
    ]


def safe_open(path):
    try:
        img = Image.open(path)
        img.load()
        return img.convert("RGB")
    except Exception as e:
        print(f"  [SKIP] {os.path.basename(path)}: {e}")
        return None


def run_predictions(model):
    results = []
    for expected_en, folder in TEST_SETS.items():
        paths = collect_images(folder)
        print(f"\n[{CLASS_KO[expected_en]}] {len(paths)}장")
        for path in paths:
            img = safe_open(path)
            if img is None:
                continue

            yolo_res = model(img, verbose=False)
            r = yolo_res[0]

            if len(r.boxes) == 0:
                pred_en = "none"
                pred_ko = "미탐지"
                conf    = 0.0
                correct = False
                ann_img = img
            else:
                dets = sorted(
                    [{"en": r.names[int(b.cls[0])], "conf": float(b.conf[0])} for b in r.boxes],
                    key=lambda x: x["conf"], reverse=True,
                )
                best    = dets[0]
                pred_en = best["en"]
                pred_ko = CLASS_KO.get(pred_en, pred_en)
                conf    = best["conf"]
                correct = (pred_en == expected_en)
                ann     = r.plot()
                ann_img = Image.fromarray(ann[..., ::-1])

            mark = "✅" if correct else "❌"
            print(f"  {os.path.basename(path):38s}  {pred_ko:10s}  {conf*100:5.1f}%  {mark}")

            results.append({
                "expected_en": expected_en,
                "expected_ko": CLASS_KO[expected_en],
                "pred_en":     pred_en,
                "pred_ko":     pred_ko,
                "conf":        conf,
                "correct":     correct,
                "path":        path,
                "original":    img,
                "annotated":   ann_img,
            })
    return results


def pil_buf(img):
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    rn = p.add_run()
    rn.text           = text
    rn.font.size      = Pt(size)
    rn.font.bold      = bold
    rn.font.color.rgb = color


def hbar(slide):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), BLUE)
    add_rect(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), BLUE)


def slide_title(prs, bangeo_n, bushiri_n):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), BLUE)
    add_text(slide, "방어 vs 부시리",
             Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "FishCheck YOLOv8 모델 — 방어·부시리 분류 정확도 테스트",
             Inches(1), Inches(3.1), Inches(11), Inches(0.8),
             size=22, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
    add_text(slide, f"방어 {bangeo_n}장  ·  부시리 {bushiri_n}장  ·  총 {bangeo_n+bushiri_n}장",
             Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             size=16, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    add_text(slide, "모델: YOLOv8s  |  데이터: data/raw/bangeo · data/raw/bushiri",
             Inches(1), Inches(4.7), Inches(11), Inches(0.5),
             size=13, color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)
    add_rect(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), BLUE)


def slide_summary(prs, results):
    slide = blank_slide(prs)
    hbar(slide)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), LGRAY)
    add_text(slide, "전체 결과 요약",
             Inches(0.4), Inches(0.12), Inches(7), Inches(0.7),
             size=26, bold=True, color=DARK)

    bangeo_r  = [r for r in results if r["expected_en"] == "bangeo"]
    bushiri_r = [r for r in results if r["expected_en"] == "bushiri"]
    total_ok  = sum(1 for r in results if r["correct"])
    total_n   = len(results)
    ba_ok     = sum(1 for r in bangeo_r  if r["correct"])
    bu_ok     = sum(1 for r in bushiri_r if r["correct"])
    total_pct = total_ok / total_n * 100 if total_n else 0
    ba_pct    = ba_ok / len(bangeo_r)  * 100 if bangeo_r  else 0
    bu_pct    = bu_ok / len(bushiri_r) * 100 if bushiri_r else 0

    for i, (label, ok, n, pct) in enumerate([
        ("전체",   total_ok, total_n,        total_pct),
        ("방어",   ba_ok,    len(bangeo_r),   ba_pct),
        ("부시리", bu_ok,    len(bushiri_r),  bu_pct),
    ]):
        lft = Inches(0.4 + i * 4.3)
        col = GREEN if pct >= 70 else ORANGE if pct >= 50 else RED
        add_rect(slide, lft, Inches(1.0), Inches(4.0), Inches(1.9), RGBColor(0xF0, 0xF4, 0xF8))
        add_text(slide, label, lft + Inches(0.15), Inches(1.08), Inches(3.7), Inches(0.45),
                 size=16, bold=True, color=DARK)
        add_text(slide, f"{ok}/{n}  ({pct:.0f}%)",
                 lft + Inches(0.15), Inches(1.5), Inches(3.7), Inches(0.6),
                 size=26, bold=True, color=col)

    add_text(slide, "혼동 행렬",
             Inches(0.4), Inches(3.05), Inches(4), Inches(0.45),
             size=16, bold=True, color=DARK)

    pred_keys    = ["bangeo", "bushiri", "gwangeo", "gajami", "none"]
    pred_display = ["방어(예측)", "부시리(예측)", "광어(예측)", "가자미(예측)", "미탐지"]
    matrix = {(e, p): 0 for e in ("bangeo", "bushiri") for p in pred_keys}
    for r in results:
        exp  = r["expected_en"]
        pred = r["pred_en"] if r["pred_en"] in pred_keys else "none"
        matrix[(exp, pred)] = matrix.get((exp, pred), 0) + 1

    cell_w = Inches(2.0)
    cell_h = Inches(0.58)
    mx_l   = Inches(0.4)
    mx_t   = Inches(3.55)

    add_rect(slide, mx_l, mx_t, Inches(1.8), cell_h, BLUE)
    add_text(slide, "실제 \\ 예측", mx_l + Inches(0.05), mx_t + Inches(0.1),
             Inches(1.7), cell_h, size=11, bold=True, color=WHITE)
    for j, pd in enumerate(pred_display):
        add_rect(slide, mx_l + Inches(1.8) + cell_w * j, mx_t, cell_w, cell_h, BLUE)
        add_text(slide, pd,
                 mx_l + Inches(1.8) + cell_w * j + Inches(0.05), mx_t + Inches(0.1),
                 cell_w - Inches(0.1), cell_h, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    for i, exp in enumerate(["bangeo", "bushiri"]):
        row_t = mx_t + cell_h * (i + 1)
        bg    = WHITE if i % 2 == 0 else LGRAY
        add_rect(slide, mx_l, row_t, Inches(1.8), cell_h, bg)
        add_text(slide, CLASS_KO[exp] + "(실제)",
                 mx_l + Inches(0.05), row_t + Inches(0.1), Inches(1.7), cell_h,
                 size=12, bold=True, color=DARK)
        for j, pred in enumerate(pred_keys):
            val     = matrix.get((exp, pred), 0)
            is_diag = (exp == pred)
            cell_bg = (RGBColor(0xD5, 0xF5, 0xE3) if (is_diag and val > 0) else
                       RGBColor(0xFD, 0xED, 0xEC) if (not is_diag and val > 0) else bg)
            add_rect(slide, mx_l + Inches(1.8) + cell_w * j, row_t, cell_w, cell_h, cell_bg)
            add_text(slide, str(val),
                     mx_l + Inches(1.8) + cell_w * j, row_t + Inches(0.1),
                     cell_w, cell_h, size=18, bold=True,
                     color=GREEN if (is_diag and val > 0) else (RED if val > 0 else DARK),
                     align=PP_ALIGN.CENTER)

    add_text(slide, "오분류 패턴",
             Inches(7.5), Inches(3.05), Inches(5.4), Inches(0.45),
             size=16, bold=True, color=DARK)
    wrong    = [r for r in results if not r["correct"]]
    patterns = {}
    for r in wrong:
        key = f"{r['expected_ko']} → {r['pred_ko']}"
        patterns[key] = patterns.get(key, 0) + 1
    patterns = dict(sorted(patterns.items(), key=lambda x: -x[1]))
    for idx, (pat, cnt) in enumerate(list(patterns.items())[:6]):
        t  = Inches(3.6) + Inches(0.62) * idx
        bg = RED if idx == 0 else ORANGE if idx == 1 else RGBColor(0xF0, 0xF0, 0xF0)
        add_rect(slide, Inches(7.5), t, Inches(5.4), Inches(0.55), bg)
        add_text(slide, f"{pat}  →  {cnt}건",
                 Inches(7.6), t + Inches(0.1), Inches(5.2), Inches(0.45),
                 size=14, bold=(idx == 0),
                 color=WHITE if idx <= 1 else DARK)


def slide_species_detail(prs, results, expected_en):
    species_ko = CLASS_KO[expected_en]
    subset     = [r for r in results if r["expected_en"] == expected_en]
    ok_cnt     = sum(1 for r in subset if r["correct"])
    total_cnt  = len(subset)
    pct        = ok_cnt / total_cnt * 100 if total_cnt else 0
    col        = GREEN if pct >= 70 else ORANGE if pct >= 50 else RED

    cols     = 5
    rows     = 5
    per_page = cols * rows
    pages    = [subset[i:i+per_page] for i in range(0, len(subset), per_page)]

    for page_idx, page_items in enumerate(pages):
        slide = blank_slide(prs)
        hbar(slide)
        add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), LGRAY)

        page_label = f" ({page_idx+1}/{len(pages)})" if len(pages) > 1 else ""
        add_text(slide, f"{species_ko} 상세 결과{page_label}",
                 Inches(0.4), Inches(0.12), Inches(8), Inches(0.7),
                 size=24, bold=True, color=DARK)
        add_text(slide, f"{ok_cnt}/{total_cnt}  ({pct:.0f}%)",
                 Inches(10.0), Inches(0.12), Inches(3.0), Inches(0.7),
                 size=24, bold=True, color=col, align=PP_ALIGN.RIGHT)

        cell_w  = Inches(2.55)
        cell_h  = Inches(1.28)
        start_l = Inches(0.2)
        start_t = Inches(1.0)

        for idx, r in enumerate(page_items):
            ci = idx % cols
            ri = idx // cols
            l  = start_l + cell_w * ci
            t  = start_t + cell_h * ri

            border_col = GREEN if r["correct"] else RED
            add_rect(slide, l, t, cell_w - Inches(0.06), cell_h - Inches(0.06), border_col)
            try:
                slide.shapes.add_picture(
                    pil_buf(r["annotated"]),
                    l + Inches(0.03), t + Inches(0.03),
                    cell_w - Inches(0.12), cell_h - Inches(0.35),
                )
            except Exception:
                pass
            lbl_bg = RGBColor(0xD5, 0xF5, 0xE3) if r["correct"] else RGBColor(0xFD, 0xED, 0xEC)
            add_rect(slide, l, t + cell_h - Inches(0.35), cell_w - Inches(0.06), Inches(0.29), lbl_bg)
            label = f"{'✅' if r['correct'] else '❌'} {r['pred_ko']} {r['conf']*100:.0f}%"
            add_text(slide, label,
                     l + Inches(0.03), t + cell_h - Inches(0.34),
                     cell_w - Inches(0.15), Inches(0.27),
                     size=9, color=GREEN if r["correct"] else RED,
                     align=PP_ALIGN.CENTER)


def slide_errors(prs, results):
    wrong  = [r for r in results if not r["correct"] and r["pred_en"] != "none"]
    no_det = [r for r in results if r["pred_en"] == "none"]

    slide = blank_slide(prs)
    hbar(slide)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), LGRAY)
    add_text(slide, f"오분류 사례  (오분류 {len(wrong)}건  |  미탐지 {len(no_det)}건)",
             Inches(0.4), Inches(0.12), Inches(11), Inches(0.7),
             size=22, bold=True, color=DARK)

    show    = wrong[:8]
    cols    = 4
    cell_w  = Inches(3.2)
    cell_h  = Inches(3.0)
    start_l = Inches(0.25)
    start_t = Inches(1.0)

    for idx, r in enumerate(show):
        ci = idx % cols
        ri = idx // cols
        l  = start_l + cell_w * ci
        t  = start_t + cell_h * ri

        add_rect(slide, l, t, cell_w - Inches(0.1), cell_h - Inches(0.1),
                 RGBColor(0xFD, 0xED, 0xEC))
        try:
            slide.shapes.add_picture(
                pil_buf(r["original"]),
                l + Inches(0.05), t + Inches(0.05),
                cell_w - Inches(0.2), cell_h - Inches(0.75),
            )
        except Exception:
            pass
        add_rect(slide, l, t + cell_h - Inches(0.68), cell_w - Inches(0.1), Inches(0.58),
                 RGBColor(0xC0, 0x39, 0x2B))
        add_text(slide,
                 f"실제: {r['expected_ko']}  →  예측: {r['pred_ko']}  ({r['conf']*100:.0f}%)",
                 l + Inches(0.06), t + cell_h - Inches(0.66),
                 cell_w - Inches(0.22), Inches(0.3),
                 size=10, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, os.path.basename(r["path"])[:26],
                 l + Inches(0.06), t + cell_h - Inches(0.36),
                 cell_w - Inches(0.22), Inches(0.28),
                 size=8, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

    if not show:
        add_text(slide, "오분류 사례 없음 🎉",
                 Inches(3), Inches(3), Inches(7), Inches(1),
                 size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


def slide_conclusion(prs, results):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), BLUE)

    bangeo_r  = [r for r in results if r["expected_en"] == "bangeo"]
    bushiri_r = [r for r in results if r["expected_en"] == "bushiri"]
    ba_ok     = sum(1 for r in bangeo_r  if r["correct"])
    bu_ok     = sum(1 for r in bushiri_r if r["correct"])
    total_ok  = ba_ok + bu_ok
    total_n   = len(results)
    ba_pct    = ba_ok / len(bangeo_r)  * 100 if bangeo_r  else 0
    bu_pct    = bu_ok / len(bushiri_r) * 100 if bushiri_r else 0

    add_text(slide, "분석 결론",
             Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    points = [
        (f"• 방어 정확도:   {ba_ok}/{len(bangeo_r)} ({ba_pct:.0f}%)",  RGBColor(0xAA, 0xDD, 0xFF)),
        (f"• 부시리 정확도: {bu_ok}/{len(bushiri_r)} ({bu_pct:.0f}%)", RGBColor(0xAA, 0xDD, 0xFF)),
        (f"• 전체 정확도:   {total_ok}/{total_n} ({total_ok/total_n*100:.0f}%)", RGBColor(0xAA, 0xDD, 0xFF)),
        ("", WHITE),
        ("• 방어·부시리는 같은 전갱이과(Carangidae) — 외형이 매우 유사", WHITE),
        ("• 주요 혼동 원인: 꼬리 색상·주둥이 형태 차이가 사진에서 불분명", WHITE),
        ("• 개선 방향: 두 어종 근접 촬영 데이터 + 라벨 수 균형 맞추기", WHITE),
    ]

    for i, (pt, col) in enumerate(points):
        add_text(slide, pt, Inches(1.2), Inches(1.5 + i * 0.68), Inches(11), Inches(0.65),
                 size=20 if pt.startswith("•") else 14, color=col)

    add_text(slide,
             "방어: 노란 가로줄 선명 · 주둥이 뾰족  |  부시리: 노란 줄 흐림 · 주둥이 둥글",
             Inches(1), Inches(6.3), Inches(11), Inches(0.5),
             size=14, color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)
    add_rect(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), BLUE)


def main():
    print("모델 로딩 중...")
    model = YOLO(MODEL_PATH)

    print("\n[예측 실행]")
    results = run_predictions(model)

    ok  = sum(1 for r in results if r["correct"])
    tot = len(results)
    print(f"\n전체 정확도: {ok}/{tot} ({ok/tot*100:.1f}%)")

    print("\nPPT 생성 중...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    bangeo_n  = sum(1 for r in results if r["expected_en"] == "bangeo")
    bushiri_n = sum(1 for r in results if r["expected_en"] == "bushiri")

    slide_title(prs, bangeo_n, bushiri_n)
    slide_summary(prs, results)
    slide_species_detail(prs, results, "bangeo")
    slide_species_detail(prs, results, "bushiri")
    slide_errors(prs, results)
    slide_conclusion(prs, results)

    prs.save(OUTPUT_PATH)
    print(f"\n완료! → {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
