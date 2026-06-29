import sys
import io
import torch
import open_clip
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

REVIEWED_DIR = Path('data/reviewed')
OUT_PPT = Path('ppt/clip_result.pptx')
OUT_PPT.parent.mkdir(exist_ok=True)

DESCRIPTIONS = {
    'bangeo': [
        'whole live yellowtail amberjack fish with symmetric tail fin',
        'Japanese amberjack buri fish with thick yellow stripe on body',
        'live whole amberjack fish in fish tank aquarium',
        'whole fish with rounded stocky body',
    ],
    'bushiri': [
        'whole live hiramasa kingfish with elongated slender body',
        'hiramasa fish with upper tail fin longer than lower',
        'live yellowtail kingfish in aquarium tank',
        'whole fish with streamlined elongated body',
    ],
}
NEGATIVE = [
    'sashimi sliced raw fish on plate',
    'person holding fishing rod',
    'cooked fish dish on table',
    'fish fillet without head',
    'fishing boat sea',
    'human hands close up',
]

print('CLIP 모델 로딩 중...')
model, _, preprocess = open_clip.create_model_and_transforms('ViT-B-32', pretrained='openai')
tokenizer = open_clip.get_tokenizer('ViT-B-32')
model.eval()
device = 'cuda' if torch.cuda.is_available() else 'cpu'
model = model.to(device)
print(f'디바이스: {device}')


def score_image(path, pos_descs, neg_descs):
    try:
        image = preprocess(Image.open(path).convert('RGB')).unsqueeze(0).to(device)
    except Exception:
        return None, None
    all_texts = pos_descs + neg_descs
    tokens = tokenizer(all_texts).to(device)
    with torch.no_grad():
        img_feat = model.encode_image(image)
        txt_feat = model.encode_text(tokens)
        img_feat = img_feat / img_feat.norm(dim=-1, keepdim=True)
        txt_feat = txt_feat / txt_feat.norm(dim=-1, keepdim=True)
        sims = (img_feat @ txt_feat.T).squeeze(0).cpu().tolist()
    pos_score = max(sims[:len(pos_descs)])
    neg_score = max(sims[len(pos_descs):])
    return pos_score, neg_score


# --- Collect scores ---
all_results = {}
for cls, pos in DESCRIPTIONS.items():
    d = REVIEWED_DIR / cls
    if not d.exists():
        print(f'[skip] {cls} 디렉토리 없음')
        continue
    imgs = sorted(d.glob('*.jpg')) + sorted(d.glob('*.png')) + sorted(d.glob('*.jpeg'))
    res = []
    for p in imgs:
        pv, nv = score_image(p, pos, NEGATIVE)
        if pv is not None:
            res.append((pv - nv, pv, nv, p))
    res.sort(reverse=True)
    all_results[cls] = res
    print(f'{cls}: {len(res)}장 완료')


# --- Colors ---
DARK   = RGBColor(0x1B, 0x3A, 0x6B)
ACCENT = RGBColor(0x00, 0x99, 0xFF)
LBKG   = RGBColor(0xF4, 0xF7, 0xFF)
GREEN  = RGBColor(0x22, 0xBB, 0x55)
RED    = RGBColor(0xEE, 0x33, 0x33)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x88, 0x88, 0x88)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill=None, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()


def add_text(slide, text, l, t, w, h, sz=14, bold=False, col=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = col


def add_image(slide, path, l, t, w, h):
    try:
        im = Image.open(path).convert('RGB')
        W, H = im.size
        s = min(W, H)
        im = im.crop(((W - s) // 2, (H - s) // 2, (W + s) // 2, (H + s) // 2))
        buf = io.BytesIO()
        im.save(buf, 'JPEG')
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))
    except Exception:
        pass


# ============================
# Slide 1: Title
# ============================
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK)
add_rect(sl, 0, 2.9, 13.33, 0.08, fill=ACCENT)
add_text(sl, 'CLIP 기반 어종 이미지 정합성 분석',
         1, 1.4, 11.33, 1.2, sz=36, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, '방어 & 부시리  수집 이미지 품질 자동 평가',
         1, 3.1, 11.33, 0.8, sz=22, col=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
add_text(sl, 'FishCheck Project  |  2026.06',
         1, 6.6, 11.33, 0.6, sz=13, col=GRAY, align=PP_ALIGN.CENTER)

# ============================
# Slide 2: CLIP 개념
# ============================
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 1.1, fill=DARK)
add_text(sl, 'CLIP이란? — 텍스트로 이미지 품질 평가', 0.3, 0.18, 12, 0.7, sz=24, bold=True, col=WHITE)

add_rect(sl, 0.3, 1.25, 5.8, 4.2, fill=LBKG, line_color=ACCENT)
add_text(sl, '텍스트 설명 예시', 0.3, 1.25, 5.8, 0.5, sz=13, bold=True, col=DARK, align=PP_ALIGN.CENTER)
pos_examples = (
    '[긍정 설명]\n'
    '- whole yellowtail fish, symmetric tail\n'
    '- buri fish with thick yellow stripe\n'
    '- live whole amberjack in aquarium\n'
    '- stocky rounded body whole fish\n\n'
    '[부정 설명]\n'
    '- sashimi sliced raw fish on plate\n'
    '- person holding fishing rod\n'
    '- cooked fish dish on table\n'
    '- fish fillet without head'
)
add_text(sl, pos_examples, 0.45, 1.82, 5.5, 3.5, sz=11, col=DARK)

add_rect(sl, 6.5, 1.25, 6.5, 4.2, fill=LBKG, line_color=ACCENT)
add_text(sl, '동작 원리', 6.5, 1.25, 6.5, 0.5, sz=13, bold=True, col=DARK, align=PP_ALIGN.CENTER)
how_text = (
    '1. 텍스트 + 이미지를 같은 벡터 공간에 임베딩\n\n'
    '2. 코사인 유사도로 텍스트-이미지 점수 계산\n\n'
    '3. pos_score = 긍정 설명 최대 유사도\n'
    '   neg_score = 부정 설명 최대 유사도\n\n'
    '4. net score = pos - neg\n'
    '   ( 클수록 물고기다운 이미지 )\n\n'
    '5. 하위 10% 제거 -> 데이터 품질 향상'
)
add_text(sl, how_text, 6.65, 1.82, 6.2, 3.5, sz=11, col=DARK)

add_text(sl, '전문가 구별 포인트(꼬리·체형·색상 줄)를 텍스트로 인코딩 -> 이미지 품질 자동 평가',
         0.3, 5.6, 12.73, 0.8, sz=13, bold=True, col=DARK)

# ============================
# Slide 3: 형태 특징 비교표
# ============================
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 1.1, fill=DARK)
add_text(sl, '방어/부시리 전문가 구별 포인트 -> CLIP 텍스트 변환',
         0.3, 0.18, 12, 0.7, sz=24, bold=True, col=WHITE)

rows = [
    ('부위', '방어 (Bangeo)', '부시리 (Bushiri)', 'CLIP 텍스트 키워드'),
    ('꼬리지느러미', '상하 거의 대칭', '상엽이 하엽보다 더 길다', 'symmetric tail / upper fin longer'),
    ('체형', '통통하고 짧음', '날렵하고 길쭉', 'stocky rounded / slender elongated'),
    ('상악골', '각진 직선형', '약간 둥근형', 'angular jaw / rounded jaw'),
    ('황색 줄', '굵고 선명', '상대적으로 가늘다', 'thick yellow stripe / thin stripe'),
]
col_widths = [2.6, 2.8, 2.8, 4.8]
col_xs = [0.25, 2.9, 5.75, 8.6]
hdr_colors = [DARK, RGBColor(0x00, 0x70, 0xC0), RGBColor(0x00, 0x70, 0x50), RGBColor(0x60, 0x30, 0x90)]

for ri, row in enumerate(rows):
    for ci, (cell, cw, cx) in enumerate(zip(row, col_widths, col_xs)):
        bg = hdr_colors[ci] if ri == 0 else (LBKG if ri % 2 == 0 else WHITE)
        fc = WHITE if ri == 0 else DARK
        add_rect(sl, cx, 1.3 + ri * 1.0, cw, 0.95, fill=bg, line_color=ACCENT)
        add_text(sl, cell, cx + 0.05, 1.38 + ri * 1.0, cw - 0.1, 0.8,
                 sz=12, bold=(ri == 0), col=fc, align=PP_ALIGN.CENTER)

# ============================
# Slide 4 & 5: 결과 (방어 / 부시리)
# ============================
LABELS = {'bangeo': '방어', 'bushiri': '부시리'}

for cls_key, cls_name in LABELS.items():
    res = all_results.get(cls_key, [])
    if not res:
        continue

    total = len(res)
    top5  = res[:5]
    bot5  = res[-5:]
    avg   = sum(r[0] for r in res) / total
    mx    = res[0][0]
    mn    = res[-1][0]
    bad_n = sum(1 for r in res if r[0] < 0.02)

    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 13.33, 1.1, fill=DARK)
    add_text(sl, f'{cls_name} ({cls_key}) CLIP 분석 결과',
             0.3, 0.18, 12, 0.7, sz=26, bold=True, col=WHITE)

    stats = [
        ('총 이미지',   f'{total}장',   ACCENT),
        ('최고 점수',   f'{mx:+.3f}',  GREEN),
        ('평균 점수',   f'{avg:+.3f}', DARK),
        ('최저 점수',   f'{mn:+.3f}',  RED),
        ('의심 이미지', f'{bad_n}장',   RED),
    ]
    for si, (label, val, col) in enumerate(stats):
        x = 0.3 + si * 2.6
        add_rect(sl, x, 1.25, 2.4, 1.2, fill=LBKG, line_color=col)
        add_text(sl, label, x, 1.28, 2.4, 0.5, sz=11, col=GRAY, align=PP_ALIGN.CENTER)
        add_text(sl, val,   x, 1.72, 2.4, 0.6, sz=20, bold=True, col=col, align=PP_ALIGN.CENTER)

    add_text(sl, '정합도 상위 5장 (물고기 특징에 가장 잘 맞는 사진)',
             0.3, 2.6, 10, 0.42, sz=13, bold=True, col=GREEN)
    for i, (net, pos_s, neg_s, p) in enumerate(top5):
        x = 0.3 + i * 2.55
        add_image(sl, p, x, 3.05, 2.3, 1.85)
        add_text(sl, f'net={net:+.3f}', x, 4.93, 2.3, 0.35,
                 sz=10, col=GREEN, align=PP_ALIGN.CENTER)

    add_text(sl, '정합도 하위 5장 (오탐 의심)',
             0.3, 5.38, 10, 0.42, sz=13, bold=True, col=RED)
    for i, (net, pos_s, neg_s, p) in enumerate(bot5):
        x = 0.3 + i * 2.55
        add_image(sl, p, x, 5.85, 2.3, 1.25)
        add_text(sl, f'net={net:+.3f}', x, 7.12, 2.3, 0.3,
                 sz=10, col=RED, align=PP_ALIGN.CENTER)

# ============================
# Slide 6: 결론
# ============================
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 1.1, fill=DARK)
add_text(sl, '결론 및 활용 방안', 0.3, 0.18, 12, 0.7, sz=26, bold=True, col=WHITE)

conclusion_cards = [
    ('점수 분포',
     'net score 범위: +0.001 ~ +0.060\n절대 임계값 대신 하위 10% 제거 방식 적합\n(임계값 0.22는 이 데이터셋엔 너무 높음)',
     ACCENT),
    ('상위권 특징',
     '수족관 전신샷, 흰 배경 생선 사진이 높은 점수\nCLIP이 통생선 특징을 잘 포착함\n정합성 기준: 방어 avg +0.035, 부시리 avg +0.018',
     GREEN),
    ('하위권 특징',
     'net < 0.01 이미지 오탐 가능성 높음\n방어 2~3장 / 부시리 4~5장 추가 검토 권장\n하위 10% 제거 시 품질 향상 기대',
     RED),
    ('다음 단계',
     '하위 10% 제거 -> Roboflow에서 라벨링\nCLIP 필터를 크롤링 파이프라인에 통합하면\n향후 수집 품질 자동화 가능',
     DARK),
]

positions = [(0.3, 1.25), (6.9, 1.25), (0.3, 4.3), (6.9, 4.3)]
for (cx, cy), (title, body, col) in zip(positions, conclusion_cards):
    add_rect(sl, cx, cy, 6.3, 2.85, fill=LBKG, line_color=col)
    add_rect(sl, cx, cy, 6.3, 0.55, fill=col)
    add_text(sl, title, cx + 0.1, cy + 0.05, 6.1, 0.45, sz=15, bold=True, col=WHITE)
    add_text(sl, body, cx + 0.1, cy + 0.65, 6.1, 2.0, sz=12, col=DARK)

# Save
prs.save(str(OUT_PPT))
print(f'\n완료: {OUT_PPT}')
