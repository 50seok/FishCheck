"""
FishCheck 프로젝트 진행 현황 PPT 생성 스크립트.

사용법:
  python tools/make_progress_ppt.py

출력: ppt/FishCheck_진행현황.pptx
"""
import os
import sys

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_PATH = os.path.join(ROOT, "ppt", "FishCheck_진행현황.pptx")

DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BLUE   = RGBColor(0x1A, 0x5C, 0x96)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x86, 0x48)
RED    = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GRAY   = RGBColor(0x77, 0x77, 0x77)
ACCENT = RGBColor(0x27, 0xAE, 0xC0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, color):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def txt(slide, text, l, t, w, h, size=18, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def header_bar(slide, title, sub=None):
    rect(slide, 0, 0, SLIDE_W, Inches(0.07), BLUE)
    rect(slide, 0, 0, SLIDE_W, Inches(1.05), LIGHT)
    txt(slide, title,
        Inches(0.4), Inches(0.12), Inches(9), Inches(0.8),
        size=28, bold=True, color=DARK)
    if sub:
        txt(slide, sub,
            Inches(10.0), Inches(0.2), Inches(3.0), Inches(0.7),
            size=14, color=GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), BLUE)


# ── 슬라이드 1: 표지 ────────────────────────────────────
def slide_cover(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK)
    rect(s, 0, 0, SLIDE_W, Inches(0.1), BLUE)
    rect(s, 0, Inches(7.4), SLIDE_W, Inches(0.1), BLUE)

    txt(s, "FishCheck",
        Inches(1), Inches(1.2), Inches(11), Inches(1.5),
        size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "수산시장 어종 사기 방지 AI — 프로젝트 진행 현황 보고",
        Inches(1), Inches(2.9), Inches(11), Inches(0.9),
        size=21, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

    rect(s, Inches(4), Inches(3.9), Inches(5.3), Inches(0.04),
         RGBColor(0x44, 0x77, 0xAA))

    lines = [
        "판별 대상: 광어 · 도다리 · 가자미 · 방어 · 부시리 · 우럭 · 개볼락 (7종)",
        "스택: Python · YOLOv8 (Ultralytics) · Streamlit · Roboflow",
        "배포: Streamlit Cloud",
    ]
    for i, line in enumerate(lines):
        txt(s, line,
            Inches(1), Inches(4.1) + Inches(0.5) * i, Inches(11), Inches(0.5),
            size=15, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

    txt(s, "2026.06",
        Inches(1), Inches(6.7), Inches(11), Inches(0.5),
        size=13, color=RGBColor(0x66, 0x66, 0x88), align=PP_ALIGN.CENTER)


# ── 슬라이드 2: 프로젝트 목표 ───────────────────────────
def slide_goal(prs):
    s = blank(prs)
    header_bar(s, "프로젝트 목표", "Why FishCheck?")

    goals = [
        ("문제",      "수산시장에서 비슷한 생선 외형을 이용한 어종 사기 발생"),
        ("해결책",    "스마트폰 사진 한 장으로 AI가 즉시 어종 판별"),
        ("판별 대상", "광어 / 도다리 / 가자미 / 방어 / 부시리 / 우럭 / 개볼락 (7종)"),
        ("배포 방식", "Streamlit Cloud — 설치 없이 브라우저에서 즉시 사용"),
        ("과제 목표", "딥러닝 이미지 분류 모델 직접 구현 · 학습 · 배포까지 완료"),
    ]

    for i, (label, desc) in enumerate(goals):
        top = Inches(1.2) + Inches(1.08) * i
        bg = LIGHT if i % 2 == 0 else WHITE
        rect(s, Inches(0.3), top, Inches(12.7), Inches(0.95), bg)
        txt(s, label, Inches(0.5), top + Inches(0.18), Inches(1.8), Inches(0.6),
            size=16, bold=True, color=BLUE)
        txt(s, desc, Inches(2.5), top + Inches(0.18), Inches(10.3), Inches(0.6),
            size=15, color=DARK)


# ── 슬라이드 3: 진행 타임라인 ───────────────────────────
def slide_timeline(prs):
    s = blank(prs)
    header_bar(s, "진행 타임라인", "2026.06")

    steps = [
        ("1단계", "프로젝트 초기화",    "CLAUDE.md 기반 폴더 구조 · requirements.txt · Streamlit 앱 뼈대",   BLUE,   "완료"),
        ("2단계", "데이터 수집",        "Naver 크롤링 + Roboflow 데이터 → 7종 504장 raw 보유",               GREEN,  "완료"),
        ("3단계", "자동 이미지 필터링", "CLIP 3단계 자동 필터 — 생물 아닌 이미지 · 저품질 자동 제거",        ACCENT, "완료"),
        ("4단계", "모델 전환",          "TensorFlow EfficientNetB0 → YOLOv8 (탐지 + 분류 동시)",             ORANGE, "완료"),
        ("5단계", "부분 학습 & 테스트", "4종 학습 → 광어/가자미 91% · 방어/부시리 69.8% 달성",               GREEN,  "완료"),
        ("6단계", "7종 통합 모델 학습", "전체 7종 재학습 + 정확도 개선 후 Streamlit Cloud 배포",              GRAY,   "진행 중"),
    ]

    for i, (step, title, desc, color, status) in enumerate(steps):
        top = Inches(1.15) + Inches(0.99) * i
        rect(s, Inches(0.25), top, Inches(0.08), Inches(0.85), color)
        rect(s, Inches(0.33), top, Inches(12.7), Inches(0.85), LIGHT if i % 2 == 0 else WHITE)

        txt(s, step, Inches(0.4), top + Inches(0.16), Inches(0.9), Inches(0.55),
            size=11, bold=True, color=color)
        txt(s, title, Inches(1.3), top + Inches(0.12), Inches(3.2), Inches(0.6),
            size=15, bold=True, color=DARK)
        txt(s, desc, Inches(4.5), top + Inches(0.16), Inches(7.3), Inches(0.55),
            size=12, color=GRAY)
        status_color = GREEN if status == "완료" else ORANGE
        txt(s, status, Inches(11.8), top + Inches(0.16), Inches(1.3), Inches(0.55),
            size=12, bold=True, color=status_color, align=PP_ALIGN.RIGHT)


# ── 슬라이드 4: 데이터 현황 ─────────────────────────────
def slide_data(prs):
    s = blank(prs)
    header_bar(s, "데이터 수집 & 현황", "raw 이미지 기준")

    txt(s, "데이터 파이프라인",
        Inches(0.4), Inches(1.15), Inches(12), Inches(0.5),
        size=16, bold=True, color=BLUE)

    pipeline = [
        ("Naver\n크롤링", "crawl_naver.py"),
        ("CLIP\n자동필터", "auto_filter.py"),
        ("선별 이미지", "selected/"),
        ("YOLOv8\n학습", "Google Colab"),
        ("Streamlit\n배포", "app.py"),
    ]

    box_w = Inches(2.0)
    start_l = Inches(0.5)
    for i, (title, sub) in enumerate(pipeline):
        l = start_l + Inches(2.55) * i
        rect(s, l, Inches(1.75), box_w, Inches(1.2), BLUE)
        txt(s, title, l + Inches(0.1), Inches(1.85), box_w - Inches(0.2), Inches(0.65),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, sub, l + Inches(0.05), Inches(2.5), box_w - Inches(0.1), Inches(0.38),
            size=9, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
        if i < 4:
            txt(s, "->", l + box_w + Inches(0.1), Inches(2.1), Inches(0.45), Inches(0.6),
                size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    txt(s, "어종별 보유 이미지 수 (raw 기준 / 목표: 100장)",
        Inches(0.4), Inches(3.15), Inches(12), Inches(0.5),
        size=16, bold=True, color=BLUE)

    species = [
        ("광어",   80), ("가자미", 80), ("우럭",   80),
        ("개볼락", 79), ("도다리", 69), ("부시리", 62), ("방어",  54),
    ]

    rect(s, Inches(0.3), Inches(3.65), Inches(12.7), Inches(0.42), BLUE)
    txt(s, "어종", Inches(0.4), Inches(3.7), Inches(1.5), Inches(0.35),
        size=12, bold=True, color=WHITE)
    txt(s, "보유 수", Inches(2.0), Inches(3.7), Inches(1.2), Inches(0.35),
        size=12, bold=True, color=WHITE)
    txt(s, "달성률 (목표 100장)", Inches(3.3), Inches(3.7), Inches(4.0), Inches(0.35),
        size=12, bold=True, color=WHITE)

    for i, (name, count) in enumerate(species):
        top = Inches(4.1) + Inches(0.43) * i
        bg = LIGHT if i % 2 == 0 else WHITE
        rect(s, Inches(0.3), top, Inches(12.7), Inches(0.4), bg)

        bar_color = GREEN if count >= 80 else (ORANGE if count >= 60 else RED)
        bar_w = Inches(6.0) * (count / 100)

        txt(s, name, Inches(0.45), top + Inches(0.08), Inches(1.4), Inches(0.3),
            size=12, bold=True, color=DARK)
        txt(s, f"{count}장", Inches(2.0), top + Inches(0.08), Inches(1.2), Inches(0.3),
            size=12, color=DARK)
        rect(s, Inches(3.3), top + Inches(0.1), bar_w, Inches(0.2), bar_color)
        txt(s, f"{count}%", Inches(3.3) + bar_w + Inches(0.1), top + Inches(0.08),
            Inches(0.8), Inches(0.3), size=10, color=GRAY)

    txt(s, "방어(54장) · 부시리(62장) 추가 수집 필요",
        Inches(0.4), Inches(7.1), Inches(12), Inches(0.3),
        size=11, bold=True, color=RED)


# ── 슬라이드 5: 모델 전환 ────────────────────────────────
def slide_model(prs):
    s = blank(prs)
    header_bar(s, "모델 아키텍처 전환", "TensorFlow -> YOLOv8")

    txt(s, "기존 (초기 계획)", Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.5),
        size=17, bold=True, color=RED)
    txt(s, "현재 채택", Inches(7.5), Inches(1.2), Inches(5.5), Inches(0.5),
        size=17, bold=True, color=GREEN)
    rect(s, Inches(6.5), Inches(1.1), Inches(0.05), Inches(5.5),
         RGBColor(0xCC, 0xCC, 0xCC))

    old_items = [
        "TensorFlow 2.x",
        "EfficientNetB0 (분류 전용)",
        "입력: 224x224 RGB",
        "출력: 7종 Softmax",
        "Hugging Face Hub 저장",
        "",
        "[X] CPU 환경에서 속도 느림",
        "[X] 바운딩 박스 없음 (어디있는지 모름)",
        "[X] TensorFlow 설치 무거움",
    ]
    for i, item in enumerate(old_items):
        color = RED if item.startswith("[X]") else DARK
        txt(s, item, Inches(0.5), Inches(1.75) + Inches(0.46) * i,
            Inches(5.5), Inches(0.42), size=13, color=color)

    new_items = [
        "Ultralytics YOLOv8s",
        "Object Detection + Classification",
        "입력: 640x640 (자동 리사이즈)",
        "출력: 바운딩박스 + 어종 + 신뢰도",
        "Roboflow 데이터 + 내장 증강",
        "",
        "[O] 실시간 탐지 가능",
        "[O] 생선 위치 바운딩 박스 시각화",
        "[O] pip install ultralytics 한 줄 설치",
    ]
    for i, item in enumerate(new_items):
        color = GREEN if item.startswith("[O]") else DARK
        txt(s, item, Inches(7.5), Inches(1.75) + Inches(0.46) * i,
            Inches(5.5), Inches(0.42), size=13, color=color)

    rect(s, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.55), LIGHT)
    txt(s, "전환 이유: YOLOv8은 '어디에 있는가(탐지) + 무엇인가(분류)'를 동시에 처리 -> 앱 사용성 향상",
        Inches(0.5), Inches(6.82), Inches(12.3), Inches(0.45),
        size=13, color=BLUE)


# ── 슬라이드 6: Roboflow & 증강 ─────────────────────────
def slide_roboflow(prs):
    s = blank(prs)
    header_bar(s, "Roboflow & YOLOv8 증강", "증강이 일어난 두 단계")

    txt(s, "증강(Augmentation)은 두 단계에서 자동 적용됩니다",
        Inches(0.4), Inches(1.15), Inches(12), Inches(0.5),
        size=16, bold=True, color=BLUE)

    stages = [
        {
            "step": "Stage 1 — Roboflow Export 시",
            "color": ACCENT,
            "items": [
                "이미지를 Roboflow에 업로드 & 라벨링",
                "Export 시 증강 배율 설정 (예: 3x)",
                "원본 1장 -> 증강 포함 여러 장 생성",
                "YOLOv8 포맷으로 다운로드",
                "결과: raw/ 폴더의 이미지들",
            ]
        },
        {
            "step": "Stage 2 — YOLOv8 학습 중 (로컬/Colab)",
            "color": GREEN,
            "items": [
                "Mosaic: 4장 합쳐서 1장으로 생성",
                "Random Flip: 좌우 반전",
                "HSV 조정: 색조 · 채도 · 밝기 랜덤 변환",
                "Random Affine: 회전 · 크기 · 이동",
                "MixUp: 두 이미지 혼합",
            ]
        },
    ]

    for i, st in enumerate(stages):
        l = Inches(0.4) + Inches(6.4) * i
        rect(s, l, Inches(1.75), Inches(6.0), Inches(4.7), LIGHT)
        rect(s, l, Inches(1.75), Inches(6.0), Inches(0.5), st["color"])
        txt(s, st["step"],
            l + Inches(0.15), Inches(1.8), Inches(5.7), Inches(0.4),
            size=14, bold=True, color=WHITE)
        for j, item in enumerate(st["items"]):
            txt(s, "* " + item,
                l + Inches(0.2), Inches(2.38) + Inches(0.76) * j,
                Inches(5.6), Inches(0.68), size=13, color=DARK)

    rect(s, Inches(0.3), Inches(6.68), Inches(12.7), Inches(0.62),
         RGBColor(0xFF, 0xF3, 0xCD))
    txt(s, "방어 54장밖에 없는 이유: Roboflow 업로드 당시 원본 소스 이미지가 적었거나 증강 배율을 낮게 설정했기 때문",
        Inches(0.5), Inches(6.75), Inches(12.1), Inches(0.5),
        size=12, color=RGBColor(0x80, 0x60, 0x00))


# ── 슬라이드 7: 테스트 결과 ─────────────────────────────
def slide_results(prs):
    s = blank(prs)
    header_bar(s, "모델 테스트 결과", "부분 학습 기준 (4종)")

    results_data = [
        {
            "pair": "광어 / 가자미 · 도다리",
            "accuracy": 91,
            "correct": 91, "total": 100,
            "color": GREEN,
            "note": "눈 방향 차이가 명확해 비교적 구별 용이",
            "errors": "일부 조명 · 각도 불량 이미지에서 혼동 발생",
        },
        {
            "pair": "방어 / 부시리",
            "accuracy": 69.8,
            "correct": 88, "total": 126,
            "color": ORANGE,
            "note": "형태적 유사도가 매우 높아 구별 어려움",
            "errors": "주요 오류: 방어->부시리 · 부시리->방어 교차 오분류 23건",
        },
    ]

    for i, r in enumerate(results_data):
        top = Inches(1.2) + Inches(2.8) * i
        rect(s, Inches(0.3), top, Inches(12.7), Inches(2.55), LIGHT)
        rect(s, Inches(0.3), top, Inches(0.1), Inches(2.55), r["color"])

        txt(s, r["pair"], Inches(0.55), top + Inches(0.15), Inches(6), Inches(0.6),
            size=20, bold=True, color=DARK)

        acc_color = GREEN if r["accuracy"] >= 85 else (ORANGE if r["accuracy"] >= 70 else RED)
        txt(s, f"{r['accuracy']}%",
            Inches(10.0), top + Inches(0.1), Inches(3.0), Inches(0.7),
            size=36, bold=True, color=acc_color, align=PP_ALIGN.RIGHT)

        txt(s, f"정답 {r['correct']} / 전체 {r['total']}장",
            Inches(0.55), top + Inches(0.8), Inches(5), Inches(0.45),
            size=13, color=GRAY)
        txt(s, "Good: " + r["note"],
            Inches(0.55), top + Inches(1.25), Inches(8), Inches(0.45),
            size=12, color=GREEN)
        txt(s, "Error: " + r["errors"],
            Inches(0.55), top + Inches(1.72), Inches(10), Inches(0.45),
            size=12, color=ORANGE)

    rect(s, Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.5), LIGHT)
    txt(s, "전체 7종 통합 학습 미완료 — 위 결과는 부분 학습(4종) 기준입니다",
        Inches(0.5), Inches(6.9), Inches(12.1), Inches(0.4),
        size=12, bold=True, color=RED)


# ── 슬라이드 8: 문제점 & 다음 단계 ─────────────────────
def slide_next(prs):
    s = blank(prs)
    header_bar(s, "현재 문제점 & 다음 단계", "개선 로드맵")

    txt(s, "현재 문제점", Inches(0.4), Inches(1.15), Inches(5.8), Inches(0.45),
        size=16, bold=True, color=RED)
    txt(s, "개선 방향", Inches(7.0), Inches(1.15), Inches(5.8), Inches(0.45),
        size=16, bold=True, color=BLUE)
    rect(s, Inches(6.5), Inches(1.1), Inches(0.05), Inches(5.2),
         RGBColor(0xDD, 0xDD, 0xDD))

    problems = [
        ("데이터 부족",      "방어 54장 · 부시리 62장 — 어종당 100장 미달"),
        ("selected/ 비어있음", "auto_filter 전체 실행 필요"),
        ("7종 미학습",       "현재 4종(방어·부시리·광어·가자미)만 학습됨"),
        ("방어-부시리 혼동", "형태 유사 — 69.8% 정확도로 과제 기준 미달 가능"),
        ("모델 미배포",      "best.pt Hugging Face Hub 업로드 안 됨"),
    ]
    for i, (title, desc) in enumerate(problems):
        top = Inches(1.72) + Inches(0.92) * i
        rect(s, Inches(0.3), top, Inches(6.0), Inches(0.8),
             RGBColor(0xFD, 0xED, 0xEC))
        txt(s, "[X] " + title, Inches(0.45), top + Inches(0.06), Inches(5.7), Inches(0.33),
            size=13, bold=True, color=RED)
        txt(s, desc, Inches(0.45), top + Inches(0.42), Inches(5.7), Inches(0.33),
            size=11, color=GRAY)

    solutions = [
        ("단기", "방어·부시리 이미지 추가 크롤링 (목표 100장)"),
        ("단기", "YOLO detection crop -> 배경 제거 -> 분류 정확도 향상"),
        ("중기", "7종 전체 통합 모델 재학습 (Google Colab)"),
        ("중기", "CLIP 제로샷으로 성능 비교 베이스라인 측정"),
        ("장기", "best.pt HuggingFace 업로드 -> Streamlit Cloud 배포"),
    ]
    for i, (timing, desc) in enumerate(solutions):
        top = Inches(1.72) + Inches(0.92) * i
        if timing == "단기":
            bg = RGBColor(0xE8, 0xF5, 0xE9)
            color = GREEN
        elif timing == "중기":
            bg = RGBColor(0xE3, 0xF2, 0xFD)
            color = BLUE
        else:
            bg = RGBColor(0xF3, 0xE5, 0xF5)
            color = RGBColor(0x7B, 0x1F, 0xA2)
        rect(s, Inches(6.7), top, Inches(6.3), Inches(0.8), bg)
        txt(s, f"[{timing}]", Inches(6.85), top + Inches(0.06), Inches(1.0), Inches(0.33),
            size=11, bold=True, color=color)
        txt(s, desc, Inches(7.9), top + Inches(0.06), Inches(4.9), Inches(0.68),
            size=12, color=DARK)


# ── 슬라이드 9: 마무리 ──────────────────────────────────
def slide_summary(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK)
    rect(s, 0, 0, SLIDE_W, Inches(0.1), BLUE)
    rect(s, 0, Inches(7.4), SLIDE_W, Inches(0.1), BLUE)

    txt(s, "진행 요약",
        Inches(1), Inches(0.8), Inches(11), Inches(0.8),
        size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    items = [
        (True,  "프로젝트 구조 · 데이터 파이프라인 · CLIP 자동 필터링 완료"),
        (True,  "TensorFlow -> YOLOv8 전환 완료"),
        (True,  "4종 부분 학습: 광어/가자미 91%  |  방어/부시리 69.8% 달성"),
        (False, "7종 통합 학습 -> 데이터 보강 후 진행 예정"),
        (False, "Streamlit Cloud 최종 배포 -> 모델 허깅페이스 업로드 후 가능"),
    ]

    for i, (done, text) in enumerate(items):
        icon = "[V]" if done else "[~]"
        color = GREEN if done else ORANGE
        txt(s, f"{icon}  {text}",
            Inches(1.5), Inches(1.85) + Inches(0.9) * i, Inches(10), Inches(0.75),
            size=17, color=color)

    txt(s, "다음 우선순위: 데이터 보강 -> 7종 통합 학습 -> HuggingFace 업로드 -> 배포",
        Inches(1), Inches(6.6), Inches(11), Inches(0.55),
        size=14, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("슬라이드 생성 중...")
    slide_cover(prs)
    slide_goal(prs)
    slide_timeline(prs)
    slide_data(prs)
    slide_model(prs)
    slide_roboflow(prs)
    slide_results(prs)
    slide_next(prs)
    slide_summary(prs)

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"완료! -> {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
