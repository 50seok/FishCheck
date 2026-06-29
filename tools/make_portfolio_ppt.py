"""FishCheck 포트폴리오 PPT — 다크 포트폴리오 스타일 (PDF 템플릿 기반)"""
from pathlib import Path
import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ────────────────────────────────────────────────────
BG      = RGBColor(0x28, 0x28, 0x28)
DGRAY   = RGBColor(0x35, 0x35, 0x35)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xAA, 0xAA, 0xAA)
DIVIDER = RGBColor(0x55, 0x55, 0x55)
ACCENT  = RGBColor(0x4A, 0x9E, 0xFF)

# ── 경로 ───────────────────────────────────────────────────────────
V3  = Path("runs/fishcheck/yolov8s_det_v3")
V2  = Path("runs/fishcheck/yolov8s_det_v2")
OUT = Path("ppt/fishcheck_portfolio.pptx")
OUT.parent.mkdir(exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H  = 13.33, 7.5

# ── matplotlib 한글 폰트 설정 ──────────────────────────────────────
_font_names = [f.name for f in fm.fontManager.ttflist]
if 'Malgun Gothic' in _font_names:
    plt.rcParams['font.family'] = 'Malgun Gothic'
plt.rcParams['axes.unicode_minus'] = False


# ── 헬퍼 ──────────────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(W), Inches(H))
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return sl

def rect(sl, l, t, w, h, fill=None, lc=None, lw=1.0):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if lc:   s.line.color.rgb = lc; s.line.width = Pt(lw)
    else:    s.line.fill.background()
    return s

def txt(sl, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, font="맑은 고딕"):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font
    return tb

def line(sl, l, t, w):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Pt(1))
    s.fill.solid(); s.fill.fore_color.rgb = DIVIDER
    s.line.fill.background()

def img(sl, path, l, t, w, h=None):
    p = Path(path)
    if not p.exists():
        return
    kw = dict(left=Inches(l), top=Inches(t), width=Inches(w))
    if h: kw['height'] = Inches(h)
    sl.shapes.add_picture(str(p), **kw)

def header(sl):
    txt(sl, "FISHCHECK", 0.5, 0.22, 4, 0.35, size=9, color=LGRAY)
    txt(sl, "2026", W-2.5, 0.22, 2.3, 0.35, size=9, color=LGRAY, align=PP_ALIGN.RIGHT)

def fig_to_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='#282828', edgecolor='none')
    buf.seek(0); plt.close(fig); return buf

def section_label(sl, num, title_kr, title_en):
    txt(sl, num, 0.5, 0.75, 2, 0.4, size=12, color=LGRAY)
    txt(sl, title_kr, 0.5, 1.15, 4.5, 1.4, size=44, bold=True, color=WHITE, font="맑은 고딕")
    txt(sl, title_en, 0.5, 2.55, 4.5, 0.45, size=14, color=LGRAY)

def top_band(sl, title, subtitle=""):
    rect(sl, 0, 0, W, 2.15, fill=BG)
    header(sl)
    txt(sl, title, 0.5, 0.65, 10, 0.72, size=20, bold=True, color=WHITE)
    if subtitle:
        txt(sl, subtitle, 0.5, 1.42, 12, 0.42, size=12, color=LGRAY)


# ══════════════════════════════════════════════════════════════════
# 1. 표지
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
txt(sl, "FISHCHECK", 0.5, 0.28, 6, 0.38, size=9, color=LGRAY)
txt(sl, "2026", W-2.5, 0.28, 2.3, 0.38, size=9, color=LGRAY, align=PP_ALIGN.RIGHT)
txt(sl, "FISHCHECK.", 0.5, 3.8, 12.3, 1.85, size=82, bold=True, color=WHITE, font="Arial Black")
txt(sl, "AI 기반 어종 판별 서비스  —  광어 vs 가자미", 0.5, 5.85, 12, 0.55, size=16, color=LGRAY)
txt(sl, "YOLOv8s  ·  Streamlit Cloud  ·  Hugging Face Hub", 0.5, 6.45, 12, 0.4, size=11, color=LGRAY)


# ══════════════════════════════════════════════════════════════════
# 2. 목차
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl)
txt(sl, "CONTENTS", 0.5, 0.9, 6, 1.5, size=62, bold=True, color=WHITE, font="Arial Black")

items = [
    ("01", "프로젝트 개요"),
    ("02", "기술 스택"),
    ("03", "데이터셋"),
    ("04", "모델 학습 과정"),
    ("05", "성능 비교"),
    ("06", "서비스 아키텍처"),
    ("07", "향후 계획"),
]
sy, step = 0.88, 0.79
for i, (num, title) in enumerate(items):
    y = sy + i * step
    line(sl, 6.5, y, 6.3)
    txt(sl, num,   6.5, y+0.1, 0.8, 0.52, size=14, color=LGRAY)
    txt(sl, title, 7.4, y+0.1, 5.4, 0.52, size=17, color=WHITE)
line(sl, 6.5, sy + len(items)*step, 6.3)


# ══════════════════════════════════════════════════════════════════
# 3. 프로젝트 개요
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl)
section_label(sl, "01", "프로젝트\n개요", "Overview")

cards = [
    ("PROBLEM", "수산시장에서 광어와 가자미는\n외형이 비슷해 소비자가\n쉽게 속는 사기 피해 발생"),
    ("SOLUTION", "스마트폰 사진 한 장으로\nYOLOv8 AI 모델이 어종을\n0.5초 이내에 판별"),
    ("SERVICE", "Streamlit 웹앱으로 배포,\n사진 업로드 즉시\n광어/가자미 결과 제공"),
]
for i, (title, body) in enumerate(cards):
    x = 4.9 + i * 2.82
    rect(sl, x, 1.0, 2.6, 5.8, fill=DGRAY)
    txt(sl, title, x+0.18, 1.18, 2.3, 0.4, size=10, bold=True, color=ACCENT)
    line(sl, x+0.18, 1.63, 2.25)
    txt(sl, body, x+0.18, 1.78, 2.3, 3.8, size=14, color=WHITE)


# ══════════════════════════════════════════════════════════════════
# 4. 기술 스택
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl)
section_label(sl, "02", "기술 스택", "Tech Stack")

categories = [
    ("AI / ML", ["YOLOv8s (Ultralytics 8.4.80)", "PyTorch", "Hugging Face Hub"]),
    ("DATA", ["Roboflow (442장 + 3× 증강)", "크롤링 + YouTube 프레임 추출", "CLIP 이미지 품질 필터"]),
    ("SERVICE", ["Streamlit", "Python 3.11", "GitHub + Streamlit Cloud"]),
    ("TOOLS", ["Google Colab (T4 GPU)", "python-pptx / matplotlib", "yt-dlp / open-clip"]),
]
positions = [(5.0, 0.9), (9.2, 0.9), (5.0, 4.1), (9.2, 4.1)]
for (x, y), (cat, items_) in zip(positions, categories):
    txt(sl, cat, x, y, 3.8, 0.4, size=10, bold=True, color=ACCENT)
    line(sl, x, y+0.45, 3.8)
    for j, item in enumerate(items_):
        txt(sl, item, x, y+0.55+j*0.48, 3.8, 0.42, size=14, color=WHITE)


# ══════════════════════════════════════════════════════════════════
# 5. 데이터셋 (막대 차트)
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl)
section_label(sl, "03", "데이터셋", "Dataset")

stats = [("442장", "원본 이미지"), ("2종", "광어 / 가자미"), ("3×", "Roboflow 증강")]
for i, (val, label) in enumerate(stats):
    y = 1.1 + i * 1.55
    txt(sl, val,   0.5, y,      4.3, 0.75, size=30, bold=True, color=WHITE)
    txt(sl, label, 0.5, y+0.68, 4.3, 0.42, size=12, color=LGRAY)
    if i < 2:
        line(sl, 0.5, y+1.18, 4.0)

fig, ax = plt.subplots(figsize=(6.2, 4.2))
fig.patch.set_facecolor('#282828'); ax.set_facecolor('#282828')
species = ['광어', '가자미']
v2_cnt  = [221, 221]
v3_cnt  = [330, 330]
x = np.arange(len(species)); bw = 0.34
b1 = ax.bar(x-bw/2, v2_cnt, bw, label='v2 원본',    color='#94A3B8', alpha=0.9)
b2 = ax.bar(x+bw/2, v3_cnt, bw, label='v3 (+증강)', color='#4A9EFF', alpha=0.9)
ax.set_ylim(0, 420); ax.set_xticks(x)
ax.set_xticklabels(species, color='white', fontsize=14)
ax.tick_params(colors='white')
for sp in ['bottom','left']: ax.spines[sp].set_color('#555555')
for sp in ['top','right']:   ax.spines[sp].set_visible(False)
ax.set_ylabel('이미지 수', color='#AAAAAA', fontsize=12)
ax.legend(facecolor='#3A3A3A', edgecolor='#555555', labelcolor='white', fontsize=12)
for b in list(b1)+list(b2):
    ax.text(b.get_x()+b.get_width()/2, b.get_height()+6,
            str(int(b.get_height())), ha='center', color='white', fontsize=12)
buf = fig_to_buf(fig)
sl.shapes.add_picture(buf, Inches(4.8), Inches(1.05), Inches(8.0), Inches(5.75))


# ══════════════════════════════════════════════════════════════════
# 6. 모델 학습 과정
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl)
section_label(sl, "04", "모델 학습\n과정", "Training Pipeline")

steps = [
    ("데이터\n수집", "크롤링 +\nYouTube 추출"),
    ("품질\n검증", "CLIP 모델\n유사도 필터"),
    ("Roboflow\n업로드", "라벨링 +\n3배 증강"),
    ("YOLOv8s\n학습", "Colab T4\n50 Epoch"),
    ("HF Hub\n배포", "best.pt 저장\n& 서비스 연동"),
]
for i, (title, body) in enumerate(steps):
    x = 4.95 + i * 1.72
    if i < len(steps)-1:
        rect(sl, x+1.3, 2.85, 0.42, 0.18, fill=LGRAY)
    rect(sl, x, 1.85, 1.25, 4.2, fill=DGRAY)
    txt(sl, str(i+1),  x+0.1, 1.95, 1.05, 0.4,  size=10, bold=True, color=ACCENT)
    txt(sl, title,     x+0.1, 2.35, 1.05, 0.85, size=12, bold=True, color=WHITE)
    line(sl, x+0.1, 3.25, 1.05)
    txt(sl, body,      x+0.1, 3.38, 1.05, 1.8,  size=11, color=LGRAY)

rect(sl, 4.95, 6.25, 8.3, 0.9, fill=DGRAY)
txt(sl, "YOLOv8s  |  입력 640×640  |  클래스 2종  |  Early Stopping 적용  |  Colab T4 GPU",
    5.1, 6.42, 8.0, 0.55, size=12, color=LGRAY)


# ══════════════════════════════════════════════════════════════════
# 7. 학습 곡선 (results.png)
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
top_band(sl, "학습 결과 — Training Curves",
         "v3 모델  |  50 Epoch  |  Early Stopping  |  최고 mAP50: 58.2% (epoch 41)")
img(sl, V3/"results.png", 0.25, 2.2, 12.83, 5.1)


# ══════════════════════════════════════════════════════════════════
# 8. 성능 비교 (막대 차트)
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl)
section_label(sl, "05", "성능 비교", "Model Performance")

fig2, ax2 = plt.subplots(figsize=(6.5, 4.2))
fig2.patch.set_facecolor('#282828'); ax2.set_facecolor('#282828')
labels = ['mAP@0.5', 'Precision', 'Recall']
v2_m   = [60.2, 43.9, 61.3]   # v2 best epoch 30
v3_m   = [58.2, 57.3, 60.7]   # v3 best epoch 41
x2 = np.arange(len(labels)); bw2 = 0.34
b3 = ax2.bar(x2-bw2/2, v2_m, bw2, label='v2',   color='#94A3B8', alpha=0.9)
b4 = ax2.bar(x2+bw2/2, v3_m, bw2, label='v3 ★', color='#4A9EFF', alpha=0.9)
ax2.set_ylim(0, 82); ax2.set_xticks(x2)
ax2.set_xticklabels(labels, color='white', fontsize=13)
ax2.tick_params(colors='white')
for sp in ['bottom','left']: ax2.spines[sp].set_color('#555555')
for sp in ['top','right']:   ax2.spines[sp].set_visible(False)
ax2.set_ylabel('성능 (%)', color='#AAAAAA', fontsize=12)
ax2.legend(facecolor='#3A3A3A', edgecolor='#555555', labelcolor='white', fontsize=12)
for b in list(b3)+list(b4):
    ax2.text(b.get_x()+b.get_width()/2, b.get_height()+1.2,
             f'{b.get_height():.1f}%', ha='center', color='white', fontsize=11)
buf2 = fig_to_buf(fig2)
sl.shapes.add_picture(buf2, Inches(4.8), Inches(1.05), Inches(8.1), Inches(5.75))

highlights = [
    ("Precision\n+13.4%p", "v2 43.9% → v3 57.3%\n오탐 대폭 감소"),
    ("신뢰도\n임계값 0.65", "낮은 신뢰도 결과\n자동 필터링"),
    ("Recall\n안정", "60.7%로\n재현율 유지"),
]
for i, (val, label) in enumerate(highlights):
    y = 1.1 + i * 1.95
    txt(sl, val,   0.5, y,      4.1, 0.9, size=18, bold=True, color=WHITE)
    txt(sl, label, 0.5, y+0.88, 4.1, 0.7, size=12, color=LGRAY)
    if i < 2:
        line(sl, 0.5, y+1.62, 4.0)


# ══════════════════════════════════════════════════════════════════
# 9. 혼동 행렬
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
top_band(sl, "혼동 행렬 — Confusion Matrix",
         "정규화된 혼동 행렬 — 광어 / 가자미 두 클래스 판별 정확도 시각화")
img(sl, V3/"confusion_matrix_normalized.png", 2.5, 2.15, 8.33, 5.15)


# ══════════════════════════════════════════════════════════════════
# 10. PR Curve + 성능 지표 카드
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
top_band(sl, "정밀도-재현율 곡선 — PR Curve",
         "탐지 임계값별 Precision / Recall 트레이드오프  |  곡선 아래 면적 = mAP50")
img(sl, V3/"BoxPR_curve.png", 0.25, 2.15, 8.8, 5.15)

cards_m = [("mAP@0.5", "58.2%"), ("Precision", "57.3%"), ("Recall", "60.7%"), ("Best Epoch", "#41")]
for i, (label, val) in enumerate(cards_m):
    y = 2.3 + i * 1.28
    rect(sl, 9.3, y, 3.75, 1.1, fill=DGRAY)
    txt(sl, val,   9.45, y+0.05, 3.4, 0.58, size=24, bold=True, color=WHITE)
    txt(sl, label, 9.45, y+0.63, 3.4, 0.35, size=11, color=LGRAY)


# ══════════════════════════════════════════════════════════════════
# 11. 서비스 아키텍처
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl)
section_label(sl, "06", "서비스\n아키텍처", "Architecture")

flow = [
    ("사용자", "스마트폰\n사진 업로드"),
    ("Streamlit\nCloud", "웹 인터페이스\n파일 수신·검증"),
    ("YOLOv8s\n모델", "HF Hub에서\n가중치 로드·추론"),
    ("결과\n표시", "어종명 +\n신뢰도 출력"),
]
for i, (title, body) in enumerate(flow):
    x = 4.95 + i * 2.1
    rect(sl, x, 2.3, 1.88, 3.2, fill=DGRAY)
    txt(sl, title, x+0.12, 2.42, 1.65, 0.78, size=13, bold=True, color=ACCENT)
    line(sl, x+0.12, 3.25, 1.65)
    txt(sl, body,  x+0.12, 3.38, 1.65, 1.75, size=12, color=WHITE)
    if i < len(flow)-1:
        txt(sl, "→", x+1.9, 3.55, 0.25, 0.55, size=18, bold=True, color=LGRAY)

rect(sl, 4.95, 5.75, 8.3, 0.9, fill=DGRAY)
txt(sl, "가중치: Hugging Face Hub  |  배포: Streamlit Cloud (GitHub master 자동 연동)  |  무료 운영",
    5.1, 5.9, 8.0, 0.55, size=12, color=LGRAY)


# ══════════════════════════════════════════════════════════════════
# 12. 향후 계획
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl)
section_label(sl, "07", "향후 계획", "Future Plans")

plans = [
    ("단기 목표", [
        "어종 확대 — 방어, 부시리 추가",
        "데이터 1,000장 이상 수집",
        "mAP50 70% 이상 달성 목표",
    ]),
    ("중장기 목표", [
        "모바일 앱 전환 (React Native)",
        "실시간 카메라 탐지",
        "가격 정보 연동 및 신고 기능",
    ]),
]
for i, (term, items_) in enumerate(plans):
    x = 4.9 + i * 4.3
    txt(sl, term, x, 1.0, 4.1, 0.42, size=11, bold=True, color=ACCENT)
    line(sl, x, 1.47, 4.1)
    for j, item in enumerate(items_):
        txt(sl, f"• {item}", x, 1.6+j*0.82, 4.1, 0.72, size=14, color=WHITE)


# ══════════════════════════════════════════════════════════════════
# 13. 마무리
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
txt(sl, "FISHCHECK", 0.5, 0.28, 6, 0.38, size=9, color=LGRAY)
txt(sl, "2026", W-2.5, 0.28, 2.3, 0.38, size=9, color=LGRAY, align=PP_ALIGN.RIGHT)
txt(sl, "THANK YOU.", 0.5, 3.8, 12.5, 1.85, size=78, bold=True, color=WHITE, font="Arial Black")
txt(sl, "GitHub: github.com/50seoks/fishcheck  |  Streamlit: fishcheck.streamlit.app",
    0.5, 6.5, 12, 0.4, size=11, color=LGRAY)


# ── 저장 ─────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"저장 완료: {OUT.resolve()}")
