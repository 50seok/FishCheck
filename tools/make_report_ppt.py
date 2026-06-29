"""FishCheck 딥러닝 산출물 PPT — v3 (차분한 색상, 적정 폰트, 이미지 경계 제한)"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ────────────────────────────────────────────────────
NAVY  = RGBColor(0x1E, 0x3A, 0x5F)
BLUE  = RGBColor(0x3B, 0x82, 0xF6)
LBLUE = RGBColor(0xDB, 0xEA, 0xFE)
INK   = RGBColor(0x1E, 0x29, 0x3B)
SLATE = RGBColor(0x64, 0x74, 0x8B)
DUST  = RGBColor(0xF1, 0xF5, 0xF9)
STEEL = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SAGE  = RGBColor(0x15, 0x80, 0x3D)
LSAGE = RGBColor(0xF0, 0xFD, 0xF4)
ROSE  = RGBColor(0xB9, 0x1C, 0x1C)
LROSE = RGBColor(0xFE, 0xF2, 0xF2)
AMBER = RGBColor(0xB4, 0x50, 0x09)
LAMBER= RGBColor(0xFF, 0xF7, 0xED)
DCODE = RGBColor(0x0F, 0x17, 0x2A)
CTEXT = RGBColor(0xA5, 0xD6, 0xFF)
BLAZE = RGBColor(0x2B, 0x56, 0x9A)

V3  = "runs/fishcheck/yolov8s_det_v3"
OUT = Path("ppt/fishcheck_report.pptx")
OUT.parent.mkdir(exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.33


def R(sl, l, t, w, h, fill=None, lc=None, lw=1.0):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if lc:
        s.line.color.rgb = lc; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def T(sl, text, l, t, w, h, sz=14, bold=False, col=INK,
      align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = col


def IMG(sl, path, l, t, w, h):
    """w·h 모두 지정 — 슬라이드 경계 초과 방지 (slight distortion 허용)"""
    p = Path(path)
    if not p.exists():
        return
    sl.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))


def hdr(sl, title, pg):
    R(sl, 0, 0, W, 0.50, fill=NAVY)
    T(sl, "FishCheck — 딥러닝 산출물 보고서",
      0.35, 0.10, 10.5, 0.32, sz=11, col=RGBColor(0xBF, 0xDB, 0xFF))
    R(sl, 12.62, 0.10, 0.45, 0.32, fill=BLUE)
    T(sl, str(pg), 12.62, 0.10, 0.45, 0.32,
      sz=12, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    R(sl, 0, 0.50, W, 0.04, fill=BLUE)
    if title:
        T(sl, title, 0.35, 0.63, 12.5, 0.55, sz=20, bold=True, col=NAVY)
        R(sl, 0.35, 1.20, 0.9, 0.04, fill=BLUE)


def mcard(sl, val, label, l, t, w=2.9, h=1.05):
    R(sl, l, t, w, h, fill=NAVY)
    T(sl, val, l, t+0.07, w, 0.55, sz=26, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    R(sl, l+0.2, t+0.66, w-0.4, 0.025, fill=RGBColor(0x3B, 0x6B, 0xB0))
    T(sl, label, l, t+0.72, w, 0.27, sz=11,
      col=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)


def fbox(sl, title, body, l, t, w, h):
    R(sl, l, t, w, h, fill=LROSE, lc=ROSE, lw=1.0)
    R(sl, l, t, 0.10, h, fill=ROSE)
    T(sl, title, l+0.20, t+0.10, w-0.28, 0.36, sz=14, bold=True, col=ROSE)
    T(sl, body,  l+0.20, t+0.52, w-0.28, h-0.60, sz=13, col=INK)


def sbox(sl, title, body, l, t, w, h):
    R(sl, l, t, w, h, fill=LSAGE, lc=SAGE, lw=1.0)
    R(sl, l, t, 0.10, h, fill=SAGE)
    T(sl, title, l+0.20, t+0.10, w-0.28, 0.36, sz=14, bold=True, col=SAGE)
    T(sl, body,  l+0.20, t+0.52, w-0.28, h-0.60, sz=13, col=INK)


def bul(sl, items, l, t, w, gap=0.52, sz=14):
    for i, item in enumerate(items):
        R(sl, l, t + i*gap + 0.10, 0.07, 0.07, fill=BLUE)
        T(sl, item, l+0.20, t + i*gap, w-0.22, gap, sz=sz, col=INK)


# ══ Slide 1: 타이틀 ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, 7.5, fill=NAVY)
R(sl, 0, 0, W, 0.07, fill=BLUE)
R(sl, 0, 6.78, W, 0.04, fill=BLUE)
R(sl, 0, 6.82, W, 0.68, fill=RGBColor(0x0A, 0x1A, 0x35))
T(sl, "FishCheck", 0.9, 0.85, 6, 0.55, sz=15, bold=True, col=BLUE)
T(sl, "AI 기반 수산시장\n어종 판별 서비스",
  0.9, 1.52, 11, 2.2, sz=38, bold=True, col=WHITE)
T(sl, "광어 vs 가자미·도다리 — 사진 한 장으로 구별하기",
  0.9, 3.90, 10.5, 0.52, sz=17, col=RGBColor(0xBF, 0xDB, 0xFF))
tags4 = ["YOLOv8s", "Streamlit Cloud", "Roboflow", "Hugging Face Hub"]
for i, tg in enumerate(tags4):
    x = 0.9 + i * 2.92
    R(sl, x, 4.60, 2.65, 0.42, fill=BLAZE)
    T(sl, tg, x, 4.60, 2.65, 0.42, sz=13, col=WHITE, align=PP_ALIGN.CENTER)
T(sl, "딥러닝 산출물  |  인공지능개발 양성과정",
  0.9, 6.93, 9, 0.38, sz=13, col=RGBColor(0x93, 0xC5, 0xFD))
T(sl, "작성자: 오영석",
  9.8, 6.93, 3.2, 0.38, sz=13, col=WHITE, align=PP_ALIGN.RIGHT)


# ══ Slide 2: 프로젝트 배경 ═══════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "프로젝트 배경 — 왜 만들었나?", 2)

R(sl, 0.35, 1.38, 6.45, 5.72, fill=DUST, lc=STEEL)
T(sl, "문제 상황", 0.55, 1.55, 6.0, 0.42, sz=15, bold=True, col=NAVY)
R(sl, 0.55, 2.02, 5.9, 0.03, fill=STEEL)
probs = [
    "광어·가자미 육안 구별이 매우 어렵다",
    "수산시장에서 가자미를 광어로\n속여 파는 사례가 빈번",
    "'좌광우도' 법칙 — 일반인이\n현장에서 적용하기 어려움",
    "사진 한 장으로 즉시 판별하는\nAI 서비스가 필요하다",
]
for i, p in enumerate(probs):
    R(sl, 0.55, 2.18 + i*1.18 + 0.12, 0.08, 0.08, fill=BLUE)
    T(sl, p, 0.80, 2.18 + i*1.18, 5.8, 1.05, sz=14, col=INK)

R(sl, 7.10, 1.38, 5.88, 5.72, fill=LBLUE, lc=BLUE, lw=1.0)
T(sl, "좌광우도  (左鰈右道)",
  7.10, 1.38, 5.88, 0.75, sz=17, bold=True, col=NAVY, align=PP_ALIGN.CENTER)
R(sl, 7.10, 2.13, 5.88, 0.03, fill=BLUE)

pairs = [
    (7.28, "광어 (Hirame)",   "눈이 왼쪽",     "입 크고 이빨 날카로움", SAGE),
    (9.92, "가자미 (Karei)",  "눈이 오른쪽",   "입 작고 체형 둥글음",   SLATE),
]
for x0, label, eye, mouth, ecol in pairs:
    R(sl, x0, 2.25, 2.55, 3.45, fill=WHITE, lc=STEEL)
    T(sl, label, x0, 2.32, 2.55, 0.45, sz=14, bold=True, col=ecol, align=PP_ALIGN.CENTER)
    R(sl, x0+0.15, 2.85, 2.25, 1.55, fill=DUST)
    T(sl, "[ 어체 이미지 ]", x0+0.15, 3.45, 2.25, 0.42,
      sz=11, col=SLATE, align=PP_ALIGN.CENTER, italic=True)
    T(sl, f"• {eye}",   x0+0.12, 4.52, 2.38, 0.38, sz=13, col=ecol, bold=True)
    T(sl, f"• {mouth}", x0+0.12, 4.92, 2.38, 0.42, sz=12, col=INK)

T(sl, "→ YOLOv8s가 이 패턴을 자동 학습",
  7.22, 6.28, 5.65, 0.42, sz=13, bold=True, col=NAVY, align=PP_ALIGN.CENTER)


# ══ Slide 3: 개발 환경 ════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "개발 환경 — 사용한 도구와 기술 스택", 3)

tools = [
    ("Python 3.10",      "개발 언어 및 런타임",                   NAVY),
    ("YOLOv8s",          "Object Detection 모델\n(Ultralytics 8.4.80)", BLUE),
    ("Roboflow",         "데이터 수집 · 라벨링 · 증강",           RGBColor(0x7C, 0x3A, 0xED)),
    ("Google Colab",     "GPU 학습 환경 (Tesla T4)",               RGBColor(0xD9, 0x7A, 0x06)),
    ("Hugging Face Hub", "학습 가중치 best.pt 저장·배포",          RGBColor(0x9A, 0x3A, 0x12)),
    ("Streamlit Cloud",  "웹 서비스 무료 배포",                    SAGE),
]
for i, (name, desc, col) in enumerate(tools):
    row, c = divmod(i, 3)
    x = 0.35 + c * 4.33
    y = 1.45 + row * 2.78
    R(sl, x, y, 4.05, 2.55, fill=col)
    T(sl, name, x+0.18, y+0.18, 3.70, 0.55, sz=18, bold=True, col=WHITE)
    R(sl, x+0.18, y+0.78, 3.40, 0.03, fill=WHITE)
    T(sl, desc,  x+0.18, y+0.90, 3.70, 1.42, sz=14, col=WHITE)


# ══ Slide 4: 흐름 로드맵 ══════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "개발 흐름 — 5단계 작업 순서", 4)

steps = [
    ("1", "데이터\n수집",   "Roboflow\n442장 + 증강"),
    ("2", "모델\n학습",     "YOLOv8s\n50 epoch"),
    ("3", "성능\n평가",     "mAP50\n54.3%"),
    ("4", "서비스\n구축",   "Streamlit\n+ HF Hub"),
    ("5", "배포 &\n테스트", "Cloud\n배포 완료"),
]
step_cols = [NAVY, BLUE, RGBColor(0x0E,0x7C,0x86), SAGE, RGBColor(0x7C,0x3A,0xED)]
for i, ((num, title, desc), col) in enumerate(zip(steps, step_cols)):
    x = 0.42 + i * 2.58
    R(sl, x+0.54, 1.48, 1.10, 1.10, fill=col)
    T(sl, num, x+0.54, 1.48, 1.10, 1.10,
      sz=30, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    R(sl, x+0.08, 2.72, 2.02, 1.20, fill=col)
    T(sl, title, x+0.08, 2.72, 2.02, 1.20,
      sz=16, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    R(sl, x+0.08, 4.04, 2.02, 0.90, fill=DUST, lc=col, lw=1.2)
    T(sl, desc, x+0.08, 4.04, 2.02, 0.90, sz=12, col=col, align=PP_ALIGN.CENTER)
    if i < 4:
        T(sl, "→", x+2.18, 3.07, 0.38, 0.65, sz=20, col=SLATE, align=PP_ALIGN.CENTER)

R(sl, 0.35, 5.22, 5.82, 0.72, fill=LROSE, lc=ROSE, lw=0.8)
T(sl, "실패 ①  방어·부시리 데이터 부족 → 2클래스로 전환",
  0.55, 5.33, 5.50, 0.50, sz=13, col=ROSE)
R(sl, 6.52, 5.22, 6.46, 0.72, fill=LROSE, lc=ROSE, lw=0.8)
T(sl, "실패 ②  배포 후 cv2 오류·오탐 → 버전 고정 + 임계값 적용",
  6.72, 5.33, 6.15, 0.50, sz=13, col=ROSE)

R(sl, 0.35, 6.12, 12.63, 0.58, fill=LBLUE, lc=BLUE, lw=0.8)
T(sl, "→ 두 번의 실패를 거쳐 현재 서비스가 완성되었습니다",
  0.60, 6.21, 12.2, 0.44, sz=14, bold=True, col=NAVY)


# ══ Slide 5: Step 1 데이터 수집 ══════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "Step 1  |  데이터 수집 — Roboflow 데이터셋 구축", 5)

mcard(sl, "442장", "원본 이미지", 0.35, 1.38, 2.85, 1.05)
mcard(sl, "2종",   "어종 클래스", 3.42, 1.38, 2.85, 1.05)
mcard(sl, "x3",    "증강 배수",   6.49, 1.38, 2.85, 1.05)

R(sl, 0.35, 2.62, 9.20, 4.48, fill=DUST, lc=STEEL)
T(sl, "데이터셋 구성", 0.55, 2.76, 9.0, 0.42, sz=15, bold=True, col=NAVY)
R(sl, 0.55, 3.24, 8.80, 0.03, fill=STEEL)
items5 = [
    "플랫폼: Roboflow  (fishcheck-jqum0 v2)",
    "라벨 형식: YOLOv8 Object Detection  (바운딩 박스)",
    "분할: Train 70%  /  Val 20%  /  Test 10%",
    "증강: Flip · Rotate ±15° · Mosaic · Brightness ±25%",
    "원칙: 통 생선(생물) 상태 — 손질·회뜨기 제외",
]
bul(sl, items5, 0.55, 3.38, 8.90, gap=0.68, sz=14)

# train_batch image — right side: l+w=9.72+3.15=12.87✓, t+h=1.55+5.58=7.13✓
T(sl, "학습 배치 샘플", 9.75, 1.38, 3.22, 0.42, sz=13, bold=True, col=NAVY)
IMG(sl, f"{V3}/train_batch0.jpg", 9.75, 1.85, 3.20, 5.18)


# ══ Slide 6: 실패 사례 1 ══════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "실패 사례 ①  |  방어·부시리 클래스 포기", 6)

R(sl, 0.35, 1.38, 12.63, 0.58, fill=ROSE)
T(sl, "수집된 이미지 수가 너무 적어 학습 불가 → 2클래스(광어·가자미)로 방향 전환",
  0.55, 1.46, 12.3, 0.42, sz=14, bold=True, col=WHITE)

fbox(sl, "방어 (Seriola quinqueradiata)",
     "수집 이미지: 47장\n수족관·수산시장 고품질 사진 구하기 어려움\n인터넷 검색 이미지 품질 불균일",
     0.35, 2.15, 4.02, 2.65)

fbox(sl, "부시리 (Seriola lalandi)",
     "수집 이미지: 45장\n방어와 외형 유사 → 라벨 혼동 위험\nYOLOv8 학습 최소 200장 이상 필요",
     4.67, 2.15, 4.02, 2.65)

R(sl, 9.00, 2.15, 3.98, 2.65, fill=LBLUE, lc=BLUE, lw=1.2)
T(sl, "결정", 9.00, 2.15, 3.98, 0.52,
  sz=15, bold=True, col=NAVY, align=PP_ALIGN.CENTER)
R(sl, 9.00, 2.67, 3.98, 0.03, fill=BLUE)
T(sl, "방어·부시리 제외\n\n→ 광어 + 가자미·도다리\n   2클래스 집중\n\n→ 클래스 수 < 데이터 품질",
  9.18, 2.80, 3.62, 1.88, sz=14, col=NAVY)

R(sl, 0.35, 5.00, 12.63, 0.52, fill=LAMBER, lc=AMBER, lw=0.8)
T(sl, "교훈:  어종 선정 전 데이터 수집 가능성을 먼저 검토해야 합니다",
  0.55, 5.08, 12.3, 0.36, sz=13, bold=True, col=AMBER)

R(sl, 0.35, 5.70, 12.63, 1.40, fill=DUST, lc=STEEL)
T(sl, "CLIP 분석 결과  (open-clip-torch · ViT-B-32 openai)",
  0.55, 5.82, 12.0, 0.38, sz=13, bold=True, col=NAVY)
T(sl, "bangeo / bushiri net_score: +0.001 ~ +0.060 (텍스트-이미지 유사도 낮음)\n"
      "→ 절대 임계값보다 '하위 10% 제거' 방식이 더 실용적임을 확인",
  0.55, 6.26, 12.1, 0.72, sz=13, col=INK)


# ══ Slide 7: Step 2 학습 설정 ════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "Step 2  |  YOLOv8s 모델 학습 설정", 7)

configs = [
    ("모델",      "YOLOv8s  (사전학습 전이학습)"),
    ("클래스",    "gwangeo  /  gajami  (nc = 2)"),
    ("입력 크기", "640 × 640"),
    ("Epochs",    "50  (Early Stopping patience = 10)"),
    ("Batch",     "16  |  GPU: Tesla T4  (Google Colab)"),
    ("Optimizer", "SGD  lr=0.01  momentum=0.937"),
]
for i, (key, val) in enumerate(configs):
    y = 1.42 + i * 0.90
    R(sl, 0.35, y, 1.90, 0.78, fill=NAVY)
    T(sl, key, 0.35, y, 1.90, 0.78, sz=13, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    R(sl, 2.38, y, 5.50, 0.78, fill=DUST, lc=STEEL)
    T(sl, val,  2.55, y+0.16, 5.25, 0.46, sz=14, col=INK)

# Code block — l+w=8.30+4.68=12.98✓, t+h=1.42+5.72=7.14✓
R(sl, 8.30, 1.42, 4.68, 5.72, fill=DCODE)
T(sl, "train.py", 8.48, 1.50, 3.0, 0.30, sz=11, col=SLATE, italic=True)
R(sl, 8.30, 1.80, 4.68, 0.03, fill=RGBColor(0x3B, 0x4A, 0x6B))
T(sl,
  "from ultralytics import YOLO\n\n"
  "model = YOLO('yolov8s.pt')\n\n"
  "model.train(\n"
  "    data='fishcheck_v2.yaml',\n"
  "    epochs=50,\n"
  "    imgsz=640,\n"
  "    batch=16,\n"
  "    patience=10,\n"
  "    device=0,\n"
  "    project='runs/fishcheck',\n"
  "    name='yolov8s_det_v3',\n"
  ")\n"
  "# 완료 → best.pt 자동 저장\n"
  "# → Hugging Face Hub 업로드",
  8.48, 1.90, 4.42, 5.14, sz=13, col=CTEXT)


# ══ Slide 8: 학습 결과 ════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "Step 2  |  학습 결과 — Loss & mAP 학습 곡선", 8)

mcard(sl, "54.3%", "mAP50",     0.35, 1.38, 3.55, 1.00)
mcard(sl, "0.631", "Precision", 4.12, 1.38, 2.80, 1.00)
mcard(sl, "50",    "Epochs",    7.14, 1.38, 2.80, 1.00)

T(sl, "학습 곡선 — Box Loss / Cls Loss / mAP50  (train · val)",
  0.35, 2.52, 12.0, 0.42, sz=14, bold=True, col=NAVY)

# results.png — l+w=0.35+12.50=12.85✓, t+h=3.00+4.15=7.15✓
IMG(sl, f"{V3}/results.png", 0.35, 3.00, 12.50, 4.15)


# ══ Slide 9: 성능 평가 ════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "Step 3  |  성능 평가 — 혼동 행렬 & PR 곡선", 9)

T(sl, "Confusion Matrix (Normalized)",
  0.35, 1.38, 6.10, 0.42, sz=14, bold=True, col=NAVY)
# confusion_matrix: l+w=0.35+5.90=6.25✓, t+h=1.85+5.05=6.90✓
IMG(sl, f"{V3}/confusion_matrix_normalized.png", 0.35, 1.85, 5.90, 5.05)

T(sl, "Box Precision-Recall Curve",
  6.70, 1.38, 6.28, 0.42, sz=14, bold=True, col=NAVY)
# BoxPR: l+w=6.70+6.05=12.75✓, t+h=1.85+5.05=6.90✓
IMG(sl, f"{V3}/BoxPR_curve.png", 6.70, 1.85, 6.05, 5.05)


# ══ Slide 10: Step 4 서비스 구축 ════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "Step 4  |  서비스 구축 — Streamlit + Hugging Face Hub", 10)

arch = [
    ("best.pt\n학습 완료",   NAVY),
    ("HF Hub\n업로드",       BLUE),
    ("Streamlit\nApp",       RGBColor(0x0E, 0x7C, 0x86)),
    ("Cloud\n자동 배포",     SAGE),
    ("사용자\n서비스",        RGBColor(0x7C, 0x3A, 0xED)),
]
for i, (label, col) in enumerate(arch):
    x = 0.35 + i * 2.20
    R(sl, x, 1.38, 1.90, 1.22, fill=col)
    T(sl, label, x, 1.38, 1.90, 1.22,
      sz=13, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        T(sl, "→", x+1.92, 1.72, 0.26, 0.55, sz=16, col=SLATE, align=PP_ALIGN.CENTER)

R(sl, 0.35, 2.80, 6.22, 4.30, fill=DUST, lc=STEEL)
T(sl, "핵심 구현", 0.55, 2.93, 6.0, 0.38, sz=15, bold=True, col=NAVY)
R(sl, 0.55, 3.38, 5.90, 0.03, fill=STEEL)
impl = [
    "@st.cache_resource — 모델 최초 1회만 로드",
    "HF Hub에서 best.pt 자동 다운로드",
    "MIME 타입 검증 (jpg / png / webp)",
    "PIL.verify() — 손상 파일 차단",
    "바운딩 박스 + 신뢰도 수치 시각화",
]
bul(sl, impl, 0.55, 3.50, 5.90, gap=0.68, sz=14)

T(sl, "검증 예측 결과 (val_batch0_pred)",
  6.82, 2.80, 6.10, 0.38, sz=13, bold=True, col=NAVY)
# val_batch: l+w=6.82+6.05=12.87✓, t+h=3.22+3.98=7.20✓
IMG(sl, f"{V3}/val_batch0_pred.jpg", 6.82, 3.22, 6.05, 3.98)


# ══ Slide 11: 실패 사례 2 ════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "실패 사례 ②  |  배포 후 발견된 문제 & 해결 과정", 11)

R(sl, 0.35, 1.38, 12.63, 0.58, fill=ROSE)
T(sl, "Streamlit Cloud 배포 직후 실제 테스트에서 발견된 2가지 문제",
  0.55, 1.47, 12.2, 0.42, sz=14, bold=True, col=WHITE)

fbox(sl, "ImportError: import cv2 실패",
     "Streamlit Cloud 배포 직후 앱 시작 불가\n"
     "ultralytics 최신 버전이 cv2 의존성 충돌",
     0.35, 2.15, 5.98, 1.90)

sbox(sl, "해결: ultralytics==8.4.80 버전 고정",
     "requirements.txt에 정확한 버전 명시\n"
     "opencv-python-headless 명시적 추가",
     0.35, 4.22, 5.98, 2.08)

fbox(sl, "우럭이 광어로 오탐 (신뢰도 54%)",
     "광어·가자미 외 어종 업로드 시\n"
     "광어(gwangeo)로 잘못 탐지됨",
     6.65, 2.15, 6.03, 1.90)

sbox(sl, "해결: conf=0.65 임계값 + 파일 검증",
     "model(img, conf=0.65) 파라미터 설정\n"
     "신뢰도 70% 미만 → 경고 메시지 표시\n"
     "비이미지 파일: MIME 타입 서버 검증",
     6.65, 4.22, 6.03, 2.08)


# ══ Slide 12: 완성 서비스 화면 ════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "Step 5  |  완성 서비스 — Streamlit 앱 화면", 12)

# Upload panel
R(sl, 0.35, 1.38, 5.98, 5.72, fill=WHITE, lc=STEEL, lw=1.2)
R(sl, 0.35, 1.38, 5.98, 0.62, fill=DUST)
T(sl, "FishCheck",
  0.55, 1.46, 5.60, 0.46, sz=16, bold=True, col=NAVY)
T(sl, "수산시장에서 생선에 속지 않도록 — 사진 한 장으로 판별합니다",
  0.45, 2.12, 5.80, 0.42, sz=12, col=INK)
R(sl, 0.45, 2.65, 5.75, 0.45, fill=RGBColor(0xFF, 0xF3, 0xCD))
T(sl, "⚠  통 생선(생물) 상태에서만 정확합니다",
  0.62, 2.72, 5.45, 0.32, sz=12, col=RGBColor(0x85, 0x60, 0x04))
R(sl, 0.45, 3.25, 5.75, 0.50, fill=DUST, lc=STEEL)
T(sl, "📁 사진 업로드     📷 카메라 촬영",
  0.65, 3.35, 5.45, 0.32, sz=13, col=SLATE)
R(sl, 0.45, 3.90, 5.75, 2.22, fill=DUST, lc=STEEL)
T(sl, "Drag and drop file here\njpg / png / webp",
  1.78, 4.62, 3.0, 0.80, sz=13, col=SLATE, align=PP_ALIGN.CENTER, italic=True)
R(sl, 2.62, 6.28, 1.55, 0.38, fill=LBLUE, lc=BLUE)
T(sl, "Browse files", 2.62, 6.31, 1.55, 0.32, sz=12, col=BLUE, align=PP_ALIGN.CENTER)

# Result panel
R(sl, 6.70, 1.38, 6.28, 5.72, fill=WHITE, lc=STEEL, lw=1.2)
R(sl, 6.70, 1.38, 6.28, 0.62, fill=DUST)
T(sl, "탐지 결과",
  6.90, 1.46, 5.90, 0.46, sz=16, bold=True, col=NAVY)
R(sl, 6.80, 2.12, 6.08, 2.48, fill=RGBColor(0x60, 0x70, 0x80))
T(sl, "[ 광어 이미지 + 바운딩 박스 ]",
  7.90, 3.18, 3.80, 0.42, sz=12, col=WHITE, align=PP_ALIGN.CENTER, italic=True)
T(sl, "gwangeo  0.87",
  6.88, 2.18, 2.55, 0.38, sz=11, bold=True, col=RGBColor(0x86, 0xEF, 0xAC))
T(sl, "판별 결과:  광어 (gwangeo)",
  6.80, 4.72, 6.08, 0.52, sz=18, bold=True, col=NAVY, align=PP_ALIGN.CENTER)
R(sl, 6.80, 5.38, 6.08, 0.48, fill=LSAGE, lc=SAGE)
T(sl, "구별 포인트: 눈이 왼쪽. 입이 크고 이빨이 날카로움.",
  6.97, 5.46, 5.75, 0.32, sz=12, col=SAGE)
T(sl, "신뢰도  87.3%",
  9.80, 6.00, 2.90, 0.38, sz=15, bold=True, col=SAGE, align=PP_ALIGN.RIGHT)


# ══ Slide 13: 성공 & 실패 요약 ════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "프로젝트 성공 & 실패 요약", 13)

R(sl, 0.35, 1.38, 6.10, 0.55, fill=SAGE)
T(sl, "✓  성공한 것들", 0.55, 1.44, 5.90, 0.44, sz=16, bold=True, col=WHITE)
wins = [
    ("YOLOv8s 탐지 성공",       "광어·가자미 실시간 탐지 + 바운딩 박스 시각화"),
    ("Streamlit Cloud 배포",     "GitHub 연동 자동 배포 — 무료 서비스 완성"),
    ("오탐 필터링 (conf=0.65)",  "타 어종 잘못 탐지 방어, 신뢰도 경고 표시"),
    ("CLIP 품질 평가 파이프라인","데이터 품질 자동 분석 워크플로우 구축"),
]
for i, (title, desc) in enumerate(wins):
    y = 2.05 + i * 1.25
    R(sl, 0.35, y, 6.10, 1.12, fill=LSAGE, lc=SAGE, lw=0.8)
    T(sl, "✓  " + title, 0.55, y+0.08, 5.80, 0.38, sz=14, bold=True, col=SAGE)
    T(sl, desc,           0.55, y+0.55, 5.80, 0.45, sz=13, col=INK)

R(sl, 6.85, 1.38, 6.13, 0.55, fill=ROSE)
T(sl, "✕  실패와 교훈", 7.05, 1.44, 5.90, 0.44, sz=16, bold=True, col=WHITE)
fails_sum = [
    ("방어·부시리 데이터 부족",  "2클래스 전환",         "어종 선정 전 수집 가능성 확인"),
    ("cv2 ImportError (배포)",   "버전 고정(8.4.80)",    "배포·개발 환경 완전 일치 필요"),
    ("우럭 오탐 (신뢰도 54%)",   "conf=0.65 적용",       "실사 테스트로 임계값 튜닝"),
]
for i, (title, sol, lesson) in enumerate(fails_sum):
    y = 2.05 + i * 1.65
    R(sl, 6.85, y, 6.13, 1.52, fill=LROSE, lc=ROSE, lw=0.8)
    T(sl, "✕  " + title,     7.05, y+0.08, 5.80, 0.36, sz=13, bold=True, col=ROSE)
    T(sl, "→ 해결: " + sol,  7.05, y+0.50, 5.80, 0.36, sz=13, col=INK)
    R(sl, 7.05, y+0.90, 5.55, 0.03, fill=STEEL)
    T(sl, "교훈: " + lesson,  7.05, y+1.05, 5.80, 0.36, sz=12, col=AMBER, italic=True)


# ══ Slide 14: 후기 & 향후 계획 ═══════════════════════════════════
sl = prs.slides.add_slide(BLANK)
hdr(sl, "후기 & 향후 계획", 14)

R(sl, 0.35, 1.38, 12.63, 0.50, fill=NAVY)
T(sl, "후기", 0.55, 1.44, 12.3, 0.40, sz=15, bold=True, col=WHITE)
R(sl, 0.35, 1.88, 12.63, 1.90, fill=DUST, lc=STEEL)
T(sl,
  "Roboflow로 데이터를 수집하고 YOLOv8s로 어종 판별 모델을 학습하여 Streamlit Cloud에 "
  "배포하는 전체 파이프라인을 직접 구축했습니다. 단순 분류(CNN)가 아닌 Object Detection "
  "으로 접근하여 바운딩 박스로 어체 위치까지 시각화했습니다. CLIP 기반 데이터 품질 평가 "
  "파이프라인 구축과 배포 후 실제 오탐 문제를 발견·수정하는 경험을 쌓았습니다.",
  0.55, 2.00, 12.3, 1.68, sz=14, col=INK)

R(sl, 0.35, 3.98, 12.63, 0.50, fill=BLUE)
T(sl, "향후 계획", 0.55, 4.04, 12.3, 0.40, sz=15, bold=True, col=WHITE)
plans = [
    ("단기",  "2클래스 전용 재학습 → mAP50 70% 이상 목표",              NAVY),
    ("중기",  "우럭·돌돔·개볼락 등 추가 어종 데이터 수집 및 클래스 확장", BLUE),
    ("장기",  "모바일 앱 연동 및 실시간 카메라 판별 성능 최적화",          RGBColor(0x0E,0x7C,0x86)),
]
for i, (term, plan, col) in enumerate(plans):
    y = 4.68 + i * 0.88
    R(sl, 0.35, y, 1.45, 0.76, fill=col)
    T(sl, term, 0.35, y, 1.45, 0.76, sz=14, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    R(sl, 1.92, y, 11.06, 0.76, fill=DUST, lc=col, lw=1.0)
    T(sl, "→  " + plan, 2.10, y+0.14, 10.75, 0.48, sz=14, col=INK)


prs.save(str(OUT))
print(f"완료: {OUT}")
