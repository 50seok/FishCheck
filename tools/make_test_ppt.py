"""
FishCheck 모델 테스트 결과 PPT 생성 스크립트.

사용법:
  python tools/make_test_ppt.py

출력: FishCheck_테스트결과.pptx
"""
import os
import sys
import io

# 프로젝트 루트를 sys.path에 추가
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)

from PIL import Image
from ultralytics import YOLO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 설정 ────────────────────────────────────────────
MODEL_PATH = os.path.join(ROOT, "models", "best.pt")
OUTPUT_PATH = os.path.join(ROOT, "ppt", "FishCheck_테스트결과.pptx")

CLASS_KO = {
    "bangeo":  "방어",
    "bushiri": "부시리",
    "gajami":  "가자미/도다리",
    "gwangeo": "광어",
}

# 어종별 테스트 이미지 2장씩
TEST_IMAGES = {
    "bangeo":  ["data/raw/bangeo/bangeo_0055.jpg",  "data/raw/bangeo/bangeo_0068.jpg"],
    "bushiri": ["data/raw/bushiri/bushiri_0031.jpg", "data/raw/bushiri/bushiri_0034.jpg"],
    "gwangeo": ["data/raw/gwangeo/gwangeo_0065.jpg", "data/raw/gwangeo/gwangeo_0066.jpg"],
    "gajami":  ["data/raw/gajami/gajami_0077.jpg",  "data/raw/gajami/gajami_0079.jpg"],
}

# 색상
GREEN  = RGBColor(0x2E, 0x86, 0x48)
RED    = RGBColor(0xC0, 0x39, 0x2B)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BLUE   = RGBColor(0x1A, 0x5C, 0x96)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def run_predictions(model: YOLO) -> list:
    results = []
    for species, paths in TEST_IMAGES.items():
        for img_path in paths:
            full = os.path.join(ROOT, img_path)
            img = Image.open(full).convert("RGB")

            yolo_results = model(img, verbose=False)
            r = yolo_results[0]

            if len(r.boxes) == 0:
                pred_en   = "none"
                pred_ko   = "미탐지"
                conf      = 0.0
                correct   = False
                annotated = img
            else:
                dets = sorted(
                    [{"class_en": r.names[int(b.cls[0])], "conf": float(b.conf[0])} for b in r.boxes],
                    key=lambda x: x["conf"], reverse=True,
                )
                best      = dets[0]
                pred_en   = best["class_en"]
                pred_ko   = CLASS_KO.get(pred_en, pred_en)
                conf      = best["conf"]
                correct   = (pred_en == species)
                ann       = r.plot()
                annotated = Image.fromarray(ann[..., ::-1])

            results.append({
                "species":    species,
                "species_ko": CLASS_KO[species],
                "img_path":   full,
                "pred_en":    pred_en,
                "pred_ko":    pred_ko,
                "conf":       conf,
                "correct":    correct,
                "annotated":  annotated,
                "original":   img,
            })
            mark = "✅" if correct else "❌"
            print(f"  {CLASS_KO[species]:8s}  예측={pred_ko:8s}  {conf*100:.1f}%  {mark}")
    return results


def pil_to_bytes(img: Image.Image) -> io.BytesIO:
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=DARK,
             align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def make_title_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), BLUE)

    add_text(slide, "FishCheck",
             Inches(1.5), Inches(1.6), Inches(10), Inches(1.3),
             font_size=60, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)

    add_text(slide, "수산시장 어종 판별 AI — YOLOv8 모델 테스트 결과",
             Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             font_size=24, color=RGBColor(0xAA, 0xCC, 0xFF),
             align=PP_ALIGN.CENTER)

    add_text(slide, "판별 대상: 방어 · 부시리 · 광어 · 가자미/도다리",
             Inches(1.5), Inches(3.9), Inches(10), Inches(0.6),
             font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC),
             align=PP_ALIGN.CENTER)

    add_text(slide, "어종당 2장 · YOLOv8s · Roboflow 데이터 증강",
             Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
             font_size=14, color=RGBColor(0x88, 0x88, 0xAA),
             align=PP_ALIGN.CENTER)

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), BLUE)


def make_summary_slide(prs, results):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), BLUE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), RGBColor(0xF0, 0xF4, 0xF8))

    add_text(slide, "테스트 결과 요약",
             Inches(0.4), Inches(0.15), Inches(7), Inches(0.75),
             font_size=28, bold=True, color=DARK)

    correct_cnt = sum(1 for r in results if r["correct"])
    total_cnt   = len(results)
    pct         = correct_cnt / total_cnt * 100

    add_text(slide, f"{correct_cnt}/{total_cnt}  ({pct:.0f}%)",
             Inches(9.5), Inches(0.15), Inches(3.4), Inches(0.75),
             font_size=28, bold=True,
             color=GREEN if pct >= 70 else RED,
             align=PP_ALIGN.RIGHT)

    headers    = ["실제 어종", "테스트 이미지", "예측 결과", "신뢰도", "정오"]
    col_lefts  = [Inches(0.3), Inches(2.1), Inches(6.8), Inches(9.5), Inches(11.5)]
    col_widths = [Inches(1.8), Inches(4.7), Inches(2.7), Inches(2.0), Inches(1.5)]
    row_h      = Inches(0.62)
    hdr_top    = Inches(1.05)

    add_rect(slide, Inches(0.2), hdr_top, Inches(12.9), row_h, BLUE)
    for h, l, w in zip(headers, col_lefts, col_widths):
        add_text(slide, h, l, hdr_top + Inches(0.12), w, row_h,
                 font_size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    for idx, r in enumerate(results):
        top = hdr_top + row_h * (idx + 1)
        bg  = RGBColor(0xFF, 0xFF, 0xFF) if idx % 2 == 0 else RGBColor(0xF5, 0xF7, 0xFA)
        add_rect(slide, Inches(0.2), top, Inches(12.9), row_h, bg)

        fname = os.path.basename(r["img_path"])
        mark  = "✅  정확" if r["correct"] else "❌  오분류"
        c_rgb = GREEN if r["correct"] else RED

        row_data = [
            (r["species_ko"], DARK, False),
            (fname, RGBColor(0x55, 0x55, 0x55), False),
            (r["pred_ko"], c_rgb, True),
            (f"{r['conf']*100:.1f}%", c_rgb, True),
            (mark, c_rgb, True),
        ]
        for (txt, color, bld), l, w in zip(row_data, col_lefts, col_widths):
            add_text(slide, txt, l, top + Inches(0.12), w, row_h,
                     font_size=13, bold=bld, color=color)

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), BLUE)


def make_result_slide(prs, r, idx, total):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), BLUE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), RGBColor(0xF0, 0xF4, 0xF8))

    mark  = "✅  정확" if r["correct"] else "❌  오분류"
    c_rgb = GREEN if r["correct"] else RED

    add_text(slide, f"테스트 {idx}/{total}  —  {r['species_ko']} (실제)",
             Inches(0.4), Inches(0.13), Inches(7), Inches(0.75),
             font_size=22, bold=True, color=DARK)

    add_text(slide, mark,
             Inches(9.8), Inches(0.13), Inches(3.1), Inches(0.75),
             font_size=22, bold=True, color=c_rgb, align=PP_ALIGN.RIGHT)

    img_h   = Inches(5.4)
    img_w   = Inches(5.8)
    img_top = Inches(1.08)

    slide.shapes.add_picture(pil_to_bytes(r["original"]),  Inches(0.3), img_top, img_w, img_h)
    slide.shapes.add_picture(pil_to_bytes(r["annotated"]), Inches(7.0), img_top, img_w, img_h)

    add_text(slide, "원본 이미지",
             Inches(0.3), img_top + img_h + Inches(0.05), img_w, Inches(0.4),
             font_size=12, color=RGBColor(0x77, 0x77, 0x77), align=PP_ALIGN.CENTER)

    add_text(slide, "YOLOv8 탐지 결과",
             Inches(7.0), img_top + img_h + Inches(0.05), img_w, Inches(0.4),
             font_size=12, color=RGBColor(0x77, 0x77, 0x77), align=PP_ALIGN.CENTER)

    box_top = img_top + img_h + Inches(0.5)
    box_bg  = RGBColor(0xE8, 0xF4, 0xFD) if r["correct"] else RGBColor(0xFD, 0xED, 0xEC)
    add_rect(slide, Inches(0.3), box_top, Inches(12.7), Inches(0.7), box_bg)

    summary = (
        f"실제: {r['species_ko']}   →   예측: {r['pred_ko']}   "
        f"  신뢰도: {r['conf']*100:.1f}%   "
        f"  ({os.path.basename(r['img_path'])})"
    )
    add_text(slide, summary,
             Inches(0.5), box_top + Inches(0.1), Inches(12.3), Inches(0.5),
             font_size=13, color=c_rgb)

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), BLUE)


def main():
    print("모델 로딩 중...")
    model = YOLO(MODEL_PATH)

    print("\n[예측 실행]")
    results = run_predictions(model)

    correct_cnt = sum(1 for r in results if r["correct"])
    print(f"\n정확도: {correct_cnt}/{len(results)} ({correct_cnt/len(results)*100:.1f}%)")

    print("\nPPT 생성 중...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    make_title_slide(prs)
    make_summary_slide(prs, results)
    for i, r in enumerate(results, start=1):
        make_result_slide(prs, r, i, len(results))

    prs.save(OUTPUT_PATH)
    print(f"\n완료! → {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
